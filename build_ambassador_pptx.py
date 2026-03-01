#!/usr/bin/env python3
"""Build GigLift Brand Ambassador Plan PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUT = os.path.join(os.path.dirname(__file__), "docs", "GigLift_Brand_Ambassador_Plan.pptx")
BG = "/Users/bthornley/.gemini/antigravity/brain/d900a025-a940-458f-9e2e-9809b998c050/ambassador_slide_bg_1772342667231.png"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(150, 150, 170)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
DARK_BG = RGBColor(15, 15, 35)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    try:
        slide.shapes.add_picture(BG, Emu(0), Emu(0), W, H)
    except Exception:
        pass


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=16, color=LIGHT, bullet_color=PURPLE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    from pptx.util import Emu as E
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(25, 25, 50)
    shape.line.color.rgb = RGBColor(60, 60, 90)
    shape.line.width = Pt(1)

    # Title
    add_text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
             title, size=16, color=accent, bold=True)

    # Items
    y = top + Inches(0.55)
    for item in items:
        add_text(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3),
                 f"• {item}", size=12, color=LIGHT)
        y += Inches(0.28)


# ═══════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "🤝 BRAND AMBASSADOR PROGRAM", size=20, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
         "Business Development Plan", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
         "Turn your best users into growth engines", size=24, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "GigLift  •  February 2026", size=16, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2: Executive Summary
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Executive Summary", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "A scalable, low-CAC growth engine", size=16, color=PURPLE)
add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
         "GigLift's Brand Ambassador Program transforms our most engaged users — working DJs, bands, solo artists, and music teachers — into advocates who drive organic sign-ups in exchange for free product access, revenue share, and exclusive perks.",
         size=18, color=LIGHT)
add_text(s, Inches(0.8), Inches(3.2), Inches(11), Inches(1.5),
         "Ambassadors promote GigLift within their local music scenes, online communities, and professional networks, creating word-of-mouth growth that scales organically across metro areas.",
         size=18, color=LIGHT)

# Key numbers
for i, (num, label) in enumerate([("500", "Ambassadors (Yr 1)"), ("12,500", "Referred Sign-ups"), ("$660K", "Revenue Potential")]):
    x = Inches(1.5) + Inches(3.8) * i
    add_text(s, x, Inches(4.8), Inches(3), Inches(0.7), num, size=42, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.5), Inches(3), Inches(0.4), label, size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 3: Ambassador Tiers
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Ambassador Tiers", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Three levels of engagement and reward", size=16, color=PURPLE)

tiers = [
    ("🥉 Starter", AMBER, [
        "3+ referral sign-ups",
        "Free Pro plan ($19/mo value)",
        "Custom referral link",
        "Ambassador badge on EPK",
    ]),
    ("🥈 Rising", CYAN, [
        "10+ sign-ups, 1 post/week",
        "Free Unlimited ($39/mo)",
        "$5 per paid referral",
        "Early feature access",
        "Co-branded social content",
    ]),
    ("🥇 Elite", PURPLE, [
        "25+ sign-ups, 2 posts/week",
        "Free Unlimited + $10/referral",
        "15% rev share on referrals",
        "Featured on landing page",
        "Product advisory input",
    ]),
]

for i, (name, color, items) in enumerate(tiers):
    x = Inches(0.8) + Inches(4) * i
    add_card(s, x, Inches(1.8), Inches(3.6), Inches(4.5), name, items, accent=color)

# ═══════════════════════════════════════
# SLIDE 4: Ambassador Types
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Ambassador Types", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Four distinct ambassador profiles for maximum reach", size=16, color=PURPLE)

types = [
    ("🎵 Scene Leaders", PURPLE, [
        "DJs/musicians with 1K+ following",
        "Instagram, TikTok, local groups",
        "Target: Performers in their area",
    ]),
    ("📚 Music Educators", CYAN, [
        "Teachers with school networks",
        "Facebook groups, teacher forums",
        "Target: Music instructors",
    ]),
    ("🎬 Content Creators", GREEN, [
        "YouTube/TikTok music creators",
        "Tutorials, demo videos",
        "Target: Musicians nationwide",
    ]),
    ("🏢 Venue Insiders", AMBER, [
        "Bookers, promoters, managers",
        "Direct referrals, events",
        "Target: Artists they work with",
    ]),
]

for i, (name, color, items) in enumerate(types):
    x = Inches(0.5) + Inches(3.1) * i
    add_card(s, x, Inches(1.8), Inches(2.9), Inches(3.5), name, items, accent=color)

# ═══════════════════════════════════════
# SLIDE 5: Referral & Tracking System
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Referral & Tracking System", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "End-to-end referral attribution and automated payouts", size=16, color=PURPLE)

# Left column: System
add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), "Technical Implementation", [
    "Custom referral links: /sign-up?ref=CODE",
    "Real-time ambassador dashboard",
    "90-day cookie attribution window",
    "Stripe Connect automated payouts",
    "Fraud detection & flagging",
], accent=CYAN)

# Right column: Metrics
add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5), "Key Metrics Tracked", [
    "Referral sign-ups (free accounts created)",
    "Conversion rate (% upgraded to paid)",
    "MRR generated per ambassador",
    "Content output (posts/month)",
    "Retention rate (90-day referred user %)",
], accent=PURPLE)

# ═══════════════════════════════════════
# SLIDE 6: Recruitment Strategy
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Recruitment Strategy", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Three phases from seed to community", size=16, color=PURPLE)

phases = [
    ("Phase 1: Seed", "Month 1-2  •  20 Ambassadors", AMBER, [
        "Identify top 50 active users",
        "Personalized email invites",
        "Post in DJ/music forums",
        "Partner with 5 music schools",
    ]),
    ("Phase 2: Scale", "Month 3-6  •  100 Ambassadors", CYAN, [
        "Public application page",
        "Ambassador referral chain",
        "Convention presence (DJ Expo)",
        "10 influencer partnerships",
    ]),
    ("Phase 3: Community", "Month 6-12  •  500 Ambassadors", PURPLE, [
        "Private Discord/Slack",
        "Regional captain program",
        "Annual ambassador summit",
        "Self-sustaining growth",
    ]),
]

for i, (title, subtitle, color, items) in enumerate(phases):
    x = Inches(0.6) + Inches(4.1) * i
    add_card(s, x, Inches(1.8), Inches(3.8), Inches(4.5), title, [subtitle, ""] + items, accent=color)

# ═══════════════════════════════════════
# SLIDE 7: Content Playbook
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Content Playbook", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Monthly content kits and weekly themes", size=16, color=PURPLE)

# Monthly kit
add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), "Monthly Content Kit", [
    "3 social post templates (carousel, story, reel)",
    "1 demo video script",
    "Branded graphics & stickers",
    "Key talking points for the month",
], accent=CYAN)

# Weekly themes
weeks = [
    ("Week 1", "Lead Discovery", "\"Found 47 new venues in one click 🔍\""),
    ("Week 2", "Outreach Results", "\"Sent 15 emails in 5 min, got 3 replies ✉️\""),
    ("Week 3", "Mode Switcher", "\"Teacher Mode found 12 schools 📚\""),
    ("Week 4", "Social Proof", "\"Booked 2 gigs from GigLift leads 🎵\""),
]

y = Inches(4.6)
for week, theme, example in weeks:
    add_text(s, Inches(0.8), y, Inches(1.5), Inches(0.3), week, size=13, color=PURPLE, bold=True)
    add_text(s, Inches(2.3), y, Inches(2), Inches(0.3), theme, size=13, color=WHITE, bold=True)
    add_text(s, Inches(4.5), y, Inches(8), Inches(0.3), example, size=12, color=MUTED)
    y += Inches(0.35)

# Platform strategy
add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.5), "Platform Strategy", [
    "Instagram: Reels & Stories (3x/week)",
    "TikTok: Short demos (2x/week)",
    "YouTube: Deep-dive tutorials (2x/month)",
    "Facebook: Teacher communities (3x/week)",
    "Reddit: Value-add comments (as needed)",
], accent=PURPLE)

# ═══════════════════════════════════════
# SLIDE 8: Incentives & Gamification
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Incentives & Gamification", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Monthly challenges and milestone rewards keep ambassadors engaged", size=16, color=PURPLE)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), "Monthly Challenges", [
    "Most sign-ups: $100 / $50 / $25 gift cards",
    "Best content (community vote): Homepage feature",
    "First to 50 referrals: Custom merch box",
    "Best conversion rate: Strategy call with founders",
], accent=CYAN)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5), "Milestone Rewards", [
    "10 sign-ups → Sticker pack + social shoutout",
    "25 sign-ups → T-shirt + priority support",
    "50 sign-ups → Hoodie + lifetime Pro access",
    "100 sign-ups → Merch bundle + equity pool + advisory seat",
], accent=PURPLE)

# ═══════════════════════════════════════
# SLIDE 9: Budget Projection
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Year 1 Budget", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Estimated $71K annual investment", size=16, color=PURPLE)

# Budget items
budget = [
    ("Free plans (100 ambassadors)", "$2,900/mo", "$34,800"),
    ("Referral payouts", "$1,500/mo", "$18,000"),
    ("Content kits & design", "$500/mo", "$6,000"),
    ("Challenge prizes", "$200/mo", "$2,400"),
    ("Merch & milestones", "$300/mo", "$3,600"),
    ("Convention presence", "$400/mo", "$4,800"),
    ("Community tools", "$100/mo", "$1,200"),
]

add_text(s, Inches(0.8), Inches(1.8), Inches(5), Inches(0.35), "LINE ITEM", size=12, color=MUTED, bold=True)
add_text(s, Inches(6), Inches(1.8), Inches(2), Inches(0.35), "MONTHLY", size=12, color=MUTED, bold=True)
add_text(s, Inches(8), Inches(1.8), Inches(2), Inches(0.35), "ANNUAL", size=12, color=MUTED, bold=True)

y = Inches(2.2)
for item, monthly, annual in budget:
    add_text(s, Inches(0.8), y, Inches(5), Inches(0.3), item, size=15, color=LIGHT)
    add_text(s, Inches(6), y, Inches(2), Inches(0.3), monthly, size=15, color=WHITE)
    add_text(s, Inches(8), y, Inches(2), Inches(0.3), annual, size=15, color=WHITE, bold=True)
    y += Inches(0.38)

# Total
add_text(s, Inches(0.8), y + Inches(0.1), Inches(5), Inches(0.35), "TOTAL", size=16, color=PURPLE, bold=True)
add_text(s, Inches(6), y + Inches(0.1), Inches(2), Inches(0.35), "$5,900/mo", size=16, color=PURPLE, bold=True)
add_text(s, Inches(8), y + Inches(0.1), Inches(2), Inches(0.35), "$70,800", size=16, color=PURPLE, bold=True)

# ═══════════════════════════════════════
# SLIDE 10: ROI Projection
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "ROI Projection", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Three scenarios: conservative, moderate, aggressive", size=16, color=PURPLE)

scenarios = [
    ("Conservative", RED, [
        "50 ambassadors",
        "8 referrals each → 400 total",
        "10% conversion → 40 paid users",
        "$10,560 annual revenue",
        "ROI: -85%",
    ]),
    ("Moderate", AMBER, [
        "150 ambassadors",
        "15 referrals each → 2,250 total",
        "15% conversion → 338 paid users",
        "$89,232 annual revenue",
        "ROI: +26%",
    ]),
    ("Aggressive", GREEN, [
        "500 ambassadors",
        "25 referrals each → 12,500 total",
        "20% conversion → 2,500 paid users",
        "$660,000 annual revenue",
        "ROI: +833%",
    ]),
]

for i, (name, color, items) in enumerate(scenarios):
    x = Inches(0.6) + Inches(4.1) * i
    add_card(s, x, Inches(1.8), Inches(3.8), Inches(4.5), name, items, accent=color)

# ═══════════════════════════════════════
# SLIDE 11: Launch Timeline
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "12-Month Launch Timeline", size=36, color=WHITE, bold=True)

timeline = [
    ("Mar 2026", "Build referral tracking, recruit 20 seed ambassadors", AMBER),
    ("Apr 2026", "Launch ambassador dashboard, first content kit", AMBER),
    ("May 2026", "Public application page, first monthly challenge", CYAN),
    ("Jun 2026", "DJ Expo presence, launch ambassador Discord", CYAN),
    ("Jul-Aug 2026", "Influencer partnerships (10), reach 100 ambassadors", CYAN),
    ("Sep 2026", "Regional captain program launches", PURPLE),
    ("Oct-Nov 2026", "Annual summit planning, self-sustaining community", PURPLE),
    ("Dec 2026", "Hit 500 ambassadors target 🎯", PURPLE),
]

y = Inches(1.5)
for date, desc, color in timeline:
    add_text(s, Inches(0.8), y, Inches(2.5), Inches(0.35), date, size=15, color=color, bold=True)
    add_text(s, Inches(3.5), y, Inches(9), Inches(0.35), desc, size=15, color=LIGHT)
    y += Inches(0.55)

# ═══════════════════════════════════════
# SLIDE 12: Next Steps
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
         "Next Steps", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
         "Immediate action items to launch the program", size=16, color=PURPLE)

steps = [
    ("1", "Build referral tracking", "Add ?ref= param handling to sign-up, store in user metadata"),
    ("2", "Create /ambassadors page", "Landing page with application form, program details, tier breakdown"),
    ("3", "Draft ambassador agreement", "Terms, brand guidelines, FTC compliance, payout schedule"),
    ("4", "Design first content kit", "March 2026 social templates, demo video scripts, branded graphics"),
    ("5", "Identify seed candidates", "Pull top 50 users by scan count, leads generated, outreach sent"),
    ("6", "Set up Stripe Connect", "Automated monthly ambassador payout infrastructure"),
]

y = Inches(1.8)
for num, title, desc in steps:
    add_text(s, Inches(0.8), y, Inches(0.6), Inches(0.5), num, size=28, color=PURPLE, bold=True)
    add_text(s, Inches(1.5), y, Inches(4), Inches(0.35), title, size=18, color=WHITE, bold=True)
    add_text(s, Inches(1.5), y + Inches(0.35), Inches(10), Inches(0.3), desc, size=14, color=MUTED)
    y += Inches(0.78)

# ═══════════════════════════════════════
# SLIDE 13: Closing
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(1), Inches(2), Inches(11), Inches(1),
         "GigLift", size=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
         "Brand Ambassador Program", size=28, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
         "Turn your best users into your best growth engine.", size=20, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "giglift.vercel.app", size=16, color=CYAN, align=PP_ALIGN.CENTER)

# Save
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
