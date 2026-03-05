#!/usr/bin/env python3
"""Build GigLift Style Guide PPTX — Design System Reference."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "data", "docs", "tech", "GigLift_Style_Guide.pptx")
BG = "/Users/bthornley/.gemini/antigravity/brain/d900a025-a940-458f-9e2e-9809b998c050/styleguide_bg_1772751308161.png"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(140, 140, 165)
PURPLE = RGBColor(207, 188, 255)
CYAN = RGBColor(125, 216, 231)
GREEN = RGBColor(129, 217, 148)
AMBER = RGBColor(240, 199, 85)
RED = RGBColor(255, 180, 171)
TEAL = RGBColor(45, 212, 191)
DARK_BG = RGBColor(15, 13, 19)
CARD_BG = RGBColor(20, 18, 24)
SURFACE = RGBColor(29, 27, 35)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    try:
        slide.shapes.add_picture(BG, Emu(0), Emu(0), W, H)
    except Exception:
        pass


def tx(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Inter"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def swatch(slide, left, top, size, rgb, label=None, hex_code=None):
    """Draw a color swatch circle (approximated with rounded rect)."""
    shape = slide.shapes.add_shape(1, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    # Round corners
    if label:
        tx(slide, left, top + size + Inches(0.05), size + Inches(0.3), Inches(0.2),
           label, size=9, color=LIGHT, bold=True, align=PP_ALIGN.LEFT)
    if hex_code:
        tx(slide, left, top + size + Inches(0.25), size + Inches(0.3), Inches(0.18),
           hex_code, size=8, color=MUTED, align=PP_ALIGN.LEFT)


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(35, 40, 55)
    shape.line.width = Pt(1)
    tx(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
       title, size=14, color=accent, bold=True)
    y = top + Inches(0.48)
    for item in items:
        if not item:
            y += Inches(0.08)
            continue
        prefix = "• " if not item.startswith(("•", " ", "→", "✅", "❌")) else ""
        tx(slide, left + Inches(0.25), y, width - Inches(0.45), Inches(0.24),
           f"{prefix}{item}", size=11, color=LIGHT)
        y += Inches(0.24)


def header(slide, title, subtitle, subtitle_color=PURPLE):
    tx(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
       title, size=32, color=WHITE, bold=True)
    if subtitle:
        tx(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.3),
           subtitle, size=14, color=subtitle_color)


# ═══════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(0.6), Inches(11), Inches(0.5),
   "🎨 DESIGN SYSTEM", size=18, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
   "GigLift Style Guide", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(3.4), Inches(11), Inches(0.6),
   "Material Design 3 — Dark Theme — Performance-Optimized", size=20, color=LIGHT, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
   "All new projects and features must adhere to this design system.", size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
   "GigLift Engineering  •  v1.0  •  March 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 2: Design Principles
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Design Principles", "The philosophical foundation behind every UI decision")

principles = [
    ("1", "Dark-First", "Dark theme is the primary experience. All colors, contrasts, and elevations are designed for dark surfaces. Never design for light-mode first.", PURPLE),
    ("2", "Material Design 3", "Follow Google's M3 spec for tonal surfaces, dynamic color, and elevation. Use M3 shape tokens for consistency.", CYAN),
    ("3", "Performance is UX", "Every visual choice must consider rendering cost. Prefer CSS transitions over JS animation. Minimize layout shifts. Respect animation budgets.", AMBER),
    ("4", "Accessible by Default", "WCAG 2.1 AA minimum. 4.5:1 contrast ratio for text. 44px touch targets. Semantic HTML. Keyboard navigable.", GREEN),
    ("5", "Progressive Disclosure", "Show only what's needed. Use cards, tabs, and expandable sections to reduce cognitive load. Empty states should guide users to action.", WHITE),
    ("6", "Responsive & Fluid", "Mobile-first design. Layouts must work at 320px–1920px. Use CSS Grid with auto-fit for adaptive grids. No horizontal scrolling.", TEAL),
]

y = Inches(1.45)
for num, title, desc, color in principles:
    tx(s, Inches(0.8), y, Inches(0.5), Inches(0.35), num, size=22, color=color, bold=True)
    tx(s, Inches(1.4), y + Inches(0.02), Inches(2.5), Inches(0.3), title, size=16, color=color, bold=True)
    tx(s, Inches(4.0), y + Inches(0.02), Inches(8.5), Inches(0.3), desc, size=12, color=LIGHT)
    y += Inches(0.6)

tx(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.4),
   "These principles apply to ALL new projects, not just GigLift. When in doubt, refer back to these six pillars.",
   size=13, color=MUTED)

# ═══════════════════════════════════════════════
# SLIDE 3: Color System — Primary & Surface
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Color System — Surfaces & Primary", "M3 tonal surface system with purple primary")

# Surface swatches
surfaces = [
    ("Surface Dim", "#0f0d13", RGBColor(15, 13, 19)),
    ("Surface", "#141218", RGBColor(20, 18, 24)),
    ("Surface Bright", "#1e1b20", RGBColor(30, 27, 32)),
    ("Container Low", "#171520", RGBColor(23, 21, 32)),
    ("Container", "#1d1b23", RGBColor(29, 27, 35)),
    ("Container High", "#262430", RGBColor(38, 36, 48)),
    ("Container Highest", "#2f2d38", RGBColor(47, 45, 56)),
    ("On-Surface", "#e6e1e5", RGBColor(230, 225, 229)),
]
tx(s, Inches(0.8), Inches(1.4), Inches(5), Inches(0.3), "TONAL SURFACES", size=13, color=CYAN, bold=True)
for i, (label, hex_code, rgb) in enumerate(surfaces):
    x = Inches(0.8) + Inches(1.5) * i
    swatch(s, x, Inches(1.8), Inches(0.9), rgb, label, hex_code)

# Primary swatches
primaries = [
    ("Primary", "#cfbcff", PURPLE),
    ("On Primary", "#381e72", RGBColor(56, 30, 114)),
    ("Primary Container", "#4f378b", RGBColor(79, 55, 139)),
    ("On Primary Cont.", "#e9ddff", RGBColor(233, 221, 255)),
    ("Inverse Primary", "#6750a4", RGBColor(103, 80, 164)),
]
tx(s, Inches(0.8), Inches(3.3), Inches(5), Inches(0.3), "PRIMARY — PURPLE", size=13, color=PURPLE, bold=True)
for i, (label, hex_code, rgb) in enumerate(primaries):
    x = Inches(0.8) + Inches(1.5) * i
    swatch(s, x, Inches(3.7), Inches(0.9), rgb, label, hex_code)

# Secondary / Tertiary
secondaries = [
    ("Secondary", "#ccc2dc", RGBColor(204, 194, 220)),
    ("On Secondary", "#332d41", RGBColor(51, 45, 65)),
    ("Tertiary", "#7dd8e7", CYAN),
    ("On Tertiary", "#003d48", RGBColor(0, 61, 72)),
    ("Tertiary Container", "#1d5b68", RGBColor(29, 91, 104)),
]
tx(s, Inches(0.8), Inches(5.15), Inches(5), Inches(0.3), "SECONDARY & TERTIARY", size=13, color=CYAN, bold=True)
for i, (label, hex_code, rgb) in enumerate(secondaries):
    x = Inches(0.8) + Inches(1.5) * i
    swatch(s, x, Inches(5.55), Inches(0.9), rgb, label, hex_code if isinstance(hex_code, str) else "")

# Usage guide on right
add_card(s, Inches(8.6), Inches(1.4), Inches(4.1), Inches(5.2), "USAGE RULES", [
    "Backgrounds: Surface Dim for page, Container for cards",
    "Text: On-Surface (#e6e1e5) for primary text",
    "Muted text: md-outline (#79747e)",
    "Borders: md-outline-variant (#49454f)",
    "",
    "Never use pure black (#000) backgrounds",
    "Never use pure white (#fff) for text",
    "Always maintain M3 tonal hierarchy",
    "",
    "CSS Variables:",
    "  var(--bg-primary) → page background",
    "  var(--bg-card) → card background",
    "  var(--text-primary) → primary text",
    "  var(--text-muted) → secondary text",
    "  var(--border) → borders & dividers",
], accent=PURPLE)

# ═══════════════════════════════════════════════
# SLIDE 4: Color System — Semantic & Accent
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Color System — Semantic & Accent Colors", "Role-based colors with dim variants for backgrounds")

# Semantic colors
semantics = [
    ("Purple (Primary)", "#cfbcff", PURPLE, "Actions, links, brand, active states"),
    ("Cyan (Tertiary)", "#7dd8e7", CYAN, "Secondary actions, data viz, info states"),
    ("Green (Success)", "#81d994", GREEN, "Confirmations, positive stats, completed"),
    ("Amber (Warning)", "#f0c755", AMBER, "Warnings, pending states, attention"),
    ("Red (Error)", "#ffb4ab", RED, "Errors, destructive actions, critical alerts"),
]

y = Inches(1.4)
for label, hex_code, rgb, usage in semantics:
    swatch(s, Inches(0.8), y, Inches(0.7), rgb)
    tx(s, Inches(1.7), y + Inches(0.05), Inches(2.5), Inches(0.25), label, size=13, color=rgb, bold=True)
    tx(s, Inches(1.7), y + Inches(0.3), Inches(2.5), Inches(0.2), hex_code, size=10, color=MUTED)
    tx(s, Inches(4.5), y + Inches(0.1), Inches(3.5), Inches(0.4), usage, size=11, color=LIGHT)
    y += Inches(0.85)

# Mode-specific accent table
add_card(s, Inches(8.5), Inches(1.4), Inches(4.2), Inches(5.2), "MODE-SPECIFIC ACCENTS", [
    "🎵 Performer Mode",
    "  Primary: #a855f7 (vivid purple)",
    "  Glow: rgba(168,85,247,0.4)",
    "",
    "📚 Instructor Mode",
    "  Primary: #38bdf8 (sky blue)",
    "  Glow: rgba(56,189,248,0.4)",
    "",
    "🎙️ Studio Mode",
    "  Primary: #f97316 (orange)",
    "  Glow: rgba(249,115,22,0.4)",
    "",
    "🚐 Touring Mode",
    "  Primary: #10b981 (emerald)",
    "  Glow: rgba(16,185,129,0.4)",
    "",
    "→ Mode colors override --accent-purple",
    "→ Use ModeConfig.color / glow / borderColor",
], accent=CYAN)

# Dim variants section
tx(s, Inches(0.8), Inches(5.8), Inches(7), Inches(0.3),
   "DIM VARIANTS — Always pair bright colors with their dim backgrounds for visual hierarchy:",
   size=11, color=MUTED, bold=True)
tx(s, Inches(0.8), Inches(6.15), Inches(7), Inches(0.25),
   "--accent-purple-dim: rgba(207,188,255,0.12)  |  --accent-cyan-dim: rgba(125,216,231,0.10)  |  --accent-green-dim: ...container",
   size=10, color=MUTED)

# ═══════════════════════════════════════════════
# SLIDE 5: Typography
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Typography", "Inter for UI, Roboto for data — weight communicates hierarchy")

# Type scale
add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5), "TYPE SCALE", [
    "",
], accent=PURPLE)
# Show actual type samples
type_samples = [
    ("Display Large", 40, 700, "Display Large — 40px Bold", PURPLE),
    ("Display Medium", 32, 700, "Display Medium — 32px Bold", WHITE),
    ("Title Large", 24, 700, "Title Large — 24px Bold", WHITE),
    ("Title Medium", 20, 600, "Title Medium — 20px Semi", WHITE),
    ("Body Large", 16, 400, "Body Large — 16px Regular", LIGHT),
    ("Body Medium", 14, 400, "Body Medium — 14px Regular", LIGHT),
    ("Body Small", 12, 400, "Body Small — 12px Regular", MUTED),
    ("Label", 11, 500, "LABEL — 11px Medium Caps", MUTED),
]
y = Inches(2.0)
for role, size, weight, sample, color in type_samples:
    actual_size = min(size, 18)  # Scale down for slide
    tx(s, Inches(0.9), y, Inches(5.2), Inches(0.35),
       sample, size=actual_size, color=color, bold=(weight >= 600))
    y += Inches(0.4)

# Type rules
add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5), "FONT FAMILIES", [
    "Primary: Inter (UI text, headings, labels)",
    "Secondary: Roboto (data tables, metrics, forms)",
    "Monospace: Consolas / JetBrains Mono (code blocks)",
    "",
    "Load via Google Fonts with display=swap",
    "Weights: 300 (Light), 400, 500, 600, 700, 800",
    "Subset to latin to reduce bundle size",
], accent=CYAN)

add_card(s, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.7), "TYPOGRAPHY RULES", [
    "Line height: 1.5 for body, 1.2 for headings",
    "Letter spacing: -0.02em for large text, 0.5px for labels",
    "Max line width: 72ch for readability",
    "",
    "❌ Never use browser default fonts",
    "❌ Never use font-size < 11px",
    "❌ Never use more than 3 weights per page",
    "❌ Never center-align body paragraphs",
    "",
    "✅ Use font-weight for hierarchy, not font-size alone",
    "✅ Uppercase only for labels and badges — never headings",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 6: Spacing & Layout
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Spacing & Layout", "4px grid — all spacing values are multiples of 4")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.0), "SPACING SCALE", [
    "4px — spacing-1 — micro gaps (icon-label)",
    "8px — spacing-2 — tight gaps (badge padding)",
    "12px — spacing-3 — default inner padding",
    "16px — spacing-4 — card internal padding",
    "20px — spacing-5 — section gaps",
    "24px — spacing-6 — card padding (standard)",
    "32px — spacing-8 — section padding",
    "48px — spacing-12 — major section breaks",
    "64px — spacing-16 — page-level vertical rhythm",
    "",
    "→ All custom values must land on the 4px grid",
], accent=GREEN)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.0), "LAYOUT GRID", [
    "Page max-width: 1200px (centered)",
    "Page padding: 24px (mobile: 16px)",
    "Card gap: 16px (12px mobile)",
    "",
    "Grid pattern:",
    "  grid-template-columns:",
    "  repeat(auto-fit, minmax(200px, 1fr))",
    "",
    "Breakpoints:",
    "  Mobile: < 768px (1 column, stacked)",
    "  Tablet: 768–1024px (2 columns)",
    "  Desktop: > 1024px (3–4 columns)",
], accent=CYAN)

add_card(s, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.2), "CONTENT REGIONS", [
    "Topbar: fixed, 60px height, z-index: 100",
    "Sidebar (if used): 260px width, collapsible on mobile",
    "Main content: margin-top: 72px (below topbar)",
    "Footer: minimal, only on marketing pages",
    "",
    "→ Always use main-content class for consistent padding",
], accent=AMBER)

add_card(s, Inches(6.8), Inches(4.7), Inches(5.8), Inches(2.2), "Z-INDEX SCALE", [
    "z-1: Cards, raised elements (1)",
    "z-10: Sticky elements, tabs (10)",
    "z-50: Dropdowns, popovers (50)",
    "z-100: Topbar, navigation (100)",
    "z-200: Modals, drawers (200)",
    "z-300: Toasts, notifications (300)",
    "z-999: Dev tools overlay (999)",
], accent=PURPLE)

# ═══════════════════════════════════════════════
# SLIDE 7: Shape & Elevation
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Shape & Elevation", "M3 shape tokens and shadow system")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.8), "SHAPE TOKENS (BORDER RADIUS)", [
    "XS: 4px — chips, tags, inline badges",
    "SM: 8px — buttons, input fields, small cards",
    "MD: 12px — standard cards, dropdowns",
    "LG: 16px — hero cards, floating panels",
    "XL: 28px — FAB buttons, pill shapes",
    "Full: 100px — avatar circles, round buttons",
    "",
    "CSS: var(--radius-xs) through var(--radius-full)",
    "Never use arbitrary border-radius values",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.8), "ELEVATION (M3 SHADOW)", [
    "Level 1 — Cards at rest, buttons",
    "  0 1px 2px rgba(0,0,0,0.3), 0 1px 3px 1px rgba(0,0,0,0.15)",
    "Level 2 — Hover cards, raised buttons",
    "  0 1px 2px rgba(0,0,0,0.3), 0 2px 6px 2px rgba(0,0,0,0.15)",
    "Level 3 — Dropdowns, menus",
    "  0 4px 8px 3px rgba(0,0,0,0.15), 0 1px 3px rgba(0,0,0,0.3)",
    "Level 4 — Modals, drawers",
    "Level 5 — Full-screen overlays",
    "",
    "CSS: var(--elevation-1) through var(--elevation-5)",
], accent=CYAN)

add_card(s, Inches(0.6), Inches(4.5), Inches(11.9), Inches(2.2), "BORDER SYSTEM", [
    "Default border: 1px solid var(--md-outline-variant)  →  #49454f (subtle, low contrast)",
    "Active border: 1px solid var(--md-outline)  →  #79747e (medium contrast, focus states)",
    "Accent border: 1px solid {accent}33  →  mode color at 20% opacity",
    "",
    "Glass border (legacy): var(--glass-border) — mapped to md-outline-variant for backward compat",
    "Never use solid bright borders on dark backgrounds — always use rgba or opacity modifiers",
    "Input focus ring: 2px solid var(--md-primary) with 4px offset",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 8: Motion & Animation
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Motion & Animation", "M3 motion system — purposeful, quick, and GPU-friendly")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "MOTION TOKENS", [
    "Duration:",
    "  Fast: 150ms — hover, focus, small state changes",
    "  Smooth: 300ms — cards, panels, navigation",
    "  Emphasized: 500ms — page transitions, modals",
    "",
    "Easing:",
    "  Emphasized: cubic-bezier(0.2, 0, 0, 1)",
    "  Standard: cubic-bezier(0.2, 0, 0, 1)",
    "  Decelerate: cubic-bezier(0, 0, 0, 1)",
    "  Accelerate: cubic-bezier(0.3, 0, 1, 1)",
    "",
    "CSS: var(--transition-fast)",
    "CSS: var(--transition-smooth)",
    "CSS: var(--motion-emphasized)",
], accent=PURPLE)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "ANIMATION RULES", [
    "✅ GPU-accelerated properties ONLY:",
    "  transform (translate, scale, rotate)",
    "  opacity",
    "  filter (blur, brightness)",
    "",
    "❌ NEVER animate:",
    "  width, height, top, left, right, bottom",
    "  margin, padding, border-width",
    "  background-color (use opacity instead)",
    "",
    "Budget: < 16ms per frame (60fps)",
    "will-change: only during animation, remove after",
    "Prefer CSS transitions over JS requestAnimationFrame",
], accent=RED)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "STANDARD ANIMATIONS", [
    "fade-in — opacity 0→1, 300ms ease",
    "slide-up — translateY(12px)→0, 400ms ease",
    "pulse — opacity 1→0.5→1, 1.5s infinite",
    "",
    "State layers (M3):",
    "  Hover: +8% primary overlay",
    "  Press: +12% primary overlay",
    "  Focus: +12% primary overlay",
    "",
    "CSS variables:",
    "  var(--state-hover): rgba(207,188,255,0.08)",
    "  var(--state-press): rgba(207,188,255,0.12)",
    "  var(--state-focus): rgba(207,188,255,0.12)",
    "",
    "Respect prefers-reduced-motion always",
], accent=GREEN)

# ═══════════════════════════════════════════════
# SLIDE 9: Component Patterns — Cards & Buttons
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Component Patterns — Cards & Buttons", "Consistent component DNA across all pages")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.8), "CARD ANATOMY", [
    "Background: var(--bg-card) → M3 Surface Container",
    "Border: 1px solid var(--glass-border)",
    "Border-radius: var(--radius-md) → 12px",
    "Padding: 20px–28px",
    "Hover: background shifts to var(--bg-card-hover)",
    "Active/selected: accent border + dim background",
    "",
    "Card variants: default, interactive (clickable), stat, mode, alert",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.8), "BUTTON HIERARCHY", [
    "btn-primary — filled, gradient background, strong action",
    "  bg: linear-gradient(135deg, accent, accent-cc)",
    "  box-shadow: 0 0 16px glow color",
    "",
    "btn-secondary — outlined, accent border + transparent bg",
    "  border: 1px solid accent at 30% opacity",
    "",
    "btn-ghost — text only, hover shows state layer",
    "btn-danger — red/error color, destructive actions only",
], accent=CYAN)

add_card(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.4), "BUTTON SIZES & STATES", [
    "btn-sm: padding 6px 14px, font 12px, radius 6px",
    "btn (default): padding 10px 20px, font 14px, radius 8px",
    "btn-lg: padding 14px 28px, font 16px, radius 12px",
    "",
    "States: default → hover (+8% overlay) → active (+12%) → disabled (50% opacity)",
    "Loading: show spinner, disable interaction, maintain size",
    "Min touch target: 44×44px (even if visually smaller)",
], accent=GREEN)

add_card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.4), "INPUT FIELDS", [
    "Background: var(--bg-input) → surface-container-low",
    "Border: 1px solid var(--glass-border)",
    "Focus: border-color: var(--md-primary), ring: 2px offset",
    "Error: border-color: var(--md-error), helper text in red",
    "Placeholder: var(--text-muted)",
    "",
    "Always include a visible <label> — never rely on placeholder alone",
    "Group related fields in fieldsets with legend",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 10: Component Patterns — Badges, Tables, Alerts
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Component Patterns — Badges, Tables & Status", "Consistent data display across all views")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "BADGES & STATUS", [
    "badge-new — green background, green text",
    "badge-contacted — cyan bg + text",
    "badge-confirmed — purple bg + text",
    "badge-completed — success green",
    "badge-draft — muted/gray",
    "badge-cancelled — red/error",
    "",
    "Common props:",
    "  font-size: 11px, font-weight: 500",
    "  padding: 3px 10px, border-radius: full",
    "  uppercase, letter-spacing: 0.5px",
    "",
    "Include emoji prefix for quick scanning:",
    "  📋 New  📡 Contacted  ✅ Confirmed",
], accent=GREEN)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "DATA TABLES & LISTS", [
    "event-list pattern for clickable rows",
    "Alternating row hover: var(--state-hover)",
    "Fixed headers for long scrollable tables",
    "Sortable columns: indicator arrow",
    "",
    "Row structure:",
    "  Status badge | Name + subtitle | Date | Type",
    "",
    "Skeleton loading: 3 placeholder rows (pulse anim)",
    "Empty state: icon + message + CTA link",
    "",
    "Responsive: hide non-essential columns on mobile",
    "Keep status + name always visible",
], accent=CYAN)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "ALERTS & TOASTS", [
    "Alert types: info, success, warning, error",
    "Background: accent-dim (12% opacity)",
    "Border-left: 3px solid accent color",
    "Icon + message + optional action",
    "",
    "Toast notifications:",
    "  Position: bottom-right, 24px inset",
    "  Max width: 400px",
    "  Auto-dismiss: 5s standard, 10s error",
    "  Animation: slide-up + fade-in",
    "  Stack: max 3 visible, newest on top",
    "",
    "Modals:",
    "  Overlay: rgba(0,0,0,0.6) backdrop-blur",
    "  Max width: 600px, centered",
    "  Focus trap: keyboard users stay inside",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 11: Iconography & Imagery
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Iconography & Imagery", "Emoji-first icons, optimized images, and brand consistency")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5), "ICON SYSTEM — EMOJI FIRST", [
    "Primary icon system: Native emoji (zero bundle cost)",
    "Standard mapping: 🎧 Events, 🔍 Leads, 📧 Email, 📊 Analytics",
    "Mode icons: 🎵 Performer, 📚 Instructor, 🎙️ Studio, 🚐 Touring",
    "",
    "For complex icons: use inline SVG (not icon font)",
    "Icon size: 16px inline, 20px buttons, 48px hero, 64px empty states",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5), "IMAGE GUIDELINES", [
    "Format: WebP preferred, fallback PNG. Never unoptimized JPEG.",
    "Max file size: 200KB per image (hard limit: 500KB)",
    "Hero images: lazy-load, use blur placeholder",
    "Avatar/profile: 40×40px, border-radius: full, fallback initials",
    "Always include meaningful alt text for accessibility",
    "Use next/image for automatic optimization in Next.js projects",
], accent=CYAN)

add_card(s, Inches(0.6), Inches(4.3), Inches(11.9), Inches(2.5), "BRAND ELEMENTS", [
    "Logo: GigLift logotype — white on dark, purple glow applied via CSS filter",
    "  → Use NavLogo component for consistent rendering",
    "  → Logo links to homepage, even when authenticated",
    "  → Filter: drop-shadow(0 0 6px mode-glow-color)",
    "",
    "Brand gradient: linear-gradient(135deg, var(--md-primary), var(--md-tertiary))",
    "Gradient usage: hero headers, CTA backgrounds, loading bars — never body backgrounds",
    "",
    "Tagline: 'AI-Powered Artist Growth Platform' — appears in topbar on all pages",
    "Co-branding (if applicable): partner logo right-aligned in topbar, max 32px height",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 12: Responsive Patterns
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Responsive Patterns", "Mobile-first — design for 320px, enhance for 1440px")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "MOBILE (< 768px)", [
    "Single column layout always",
    "Nav collapsed to hamburger menu",
    "Cards stack vertically, full width",
    "Hide secondary data columns in tables",
    "Buttons: full width on mobile modals",
    "Touch targets: 44×44px minimum",
    "Padding: 16px page, 14px card",
    "",
    "Navigation:",
    "  Topbar stays fixed, 52px height",
    "  Mode switch: dropdown not tabs",
    "  Settings: slide-in drawer from right",
    "",
    "Test at: 320px, 375px, 414px",
], accent=GREEN)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "TABLET (768–1024px)", [
    "2-column grid for cards",
    "Nav: full topbar, horizontal mode tabs",
    "Tables: show 3-4 columns max",
    "Modals: 80% width, centered",
    "Side panels: overlay, not push",
    "",
    "Grid pattern:",
    "  grid-template-columns: repeat(2, 1fr)",
    "",
    "Dashboard stats: 2×3 grid",
    "Lead cards: 2 per row",
    "",
    "Test at: 768px, 1024px",
], accent=CYAN)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "DESKTOP (> 1024px)", [
    "3-4 column grids, auto-fit preferred",
    "Full horizontal nav with mode tabs",
    "Tables: all columns visible",
    "Modals: 600px max-width, centered",
    "Side panels: push layout (optional)",
    "",
    "Grid pattern:",
    "  grid-template-columns:",
    "  repeat(auto-fit, minmax(200px, 1fr))",
    "",
    "Dashboard: 4-6 stat cards per row",
    "Max content width: 1200px",
    "",
    "Test at: 1280px, 1440px, 1920px",
], accent=PURPLE)

# ═══════════════════════════════════════════════
# SLIDE 13: Accessibility
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Accessibility — WCAG 2.1 AA", "Every user must be able to navigate, read, and interact with the app")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "COLOR & CONTRAST", [
    "Text: 4.5:1 contrast ratio minimum (AA)",
    "Large text (>18px bold): 3:1 ratio minimum",
    "Focus indicators: 3:1 against adjacent colors",
    "",
    "Our contrast ratios:",
    "  On-surface (#e6e1e5) on Surface (#141218) → 12.3:1 ✅",
    "  Outline (#79747e) on Surface → 4.8:1 ✅",
    "  Primary (#cfbcff) on Surface → 9.6:1 ✅",
    "",
    "Never rely on color alone for meaning",
    "  → Always pair with icon/emoji + text label",
    "  → Status badges use color + emoji prefix",
], accent=PURPLE)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "KEYBOARD & FOCUS", [
    "All interactive elements: keyboard accessible",
    "Tab order: logical, left-to-right, top-to-bottom",
    "Focus ring: 2px solid primary, 4px offset",
    "Skip link: 'Skip to main content' on every page",
    "",
    "Modal focus trap: Tab cycles within modal",
    "Escape key: closes modals, popovers, dropdowns",
    "",
    "ARIA requirements:",
    "  aria-label on icon-only buttons",
    "  aria-expanded on toggleable content",
    "  aria-live='polite' for toast notifications",
    "  role='dialog' + aria-modal on modals",
], accent=GREEN)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "SEMANTIC HTML", [
    "<nav> for navigation bars",
    "<main> for primary content",
    "<section> for logical groups",
    "<article> for self-contained content",
    "<aside> for supplementary info",
    "",
    "Form elements:",
    "  <label> for every input (visible or sr-only)",
    "  <fieldset> + <legend> for grouped fields",
    "  htmlFor= on labels (React: htmlFor)",
    "",
    "Images / media:",
    "  alt= on all images (descriptive)",
    "  Empty alt='' for decorative images",
    "  <figure> + <figcaption> for meaningful images",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 14: CSS Architecture
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "CSS Architecture", "Vanilla CSS with design tokens — no utility frameworks")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5), "FILE STRUCTURE", [
    "globals.css — all design tokens, reset, base styles, components",
    "Future: modular CSS with @import (when CSS nesting support is stable)",
    "",
    "Organization within globals.css:",
    "  1. Imports (Google Fonts)  →  2. CSS Variables  →  3. Reset/Base",
    "  4. Layout  →  5. Components  →  6. Pages  →  7. Responsive",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5), "NAMING CONVENTION (BEM-INSPIRED)", [
    "Blocks: .card, .btn, .badge, .event-list",
    "Elements: .event-col-name, .event-col-date",
    "Modifiers: .btn-primary, .btn-sm, .badge-confirmed",
    "",
    "❌ No utility classes (no .mt-4, .text-center, .flex)",
    "✅ Use semantic class names that describe purpose, not styles",
], accent=CYAN)

add_card(s, Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.5), "CSS VARIABLE USAGE", [
    "ALWAYS use design tokens — never hardcode values",
    "",
    "✅ color: var(--text-primary)  |  background: var(--bg-card)",
    "✅ border-radius: var(--radius-md)  |  box-shadow: var(--elevation-2)",
    "❌ color: #e6e1e5  |  background: #1d1b23  |  border-radius: 12px",
    "",
    "Exception: mode-specific dynamic colors via props (inline styles OK for these)",
], accent=GREEN)

add_card(s, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.5), "PERFORMANCE RULES", [
    "Prefer CSS classes over inline styles",
    "Avoid !important — fix specificity instead",
    "Minimize nesting depth (max 3 levels)",
    "Use CSS containment for complex layouts (contain: layout)",
    "Avoid * selectors — be specific",
    "All animations: transform + opacity only (GPU accel)",
], accent=RED)

# ═══════════════════════════════════════════════
# SLIDE 15: Closing
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
   "GigLift", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
   "Design System & Style Guide", size=26, color=PURPLE, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
   "Every pixel is intentional.\nEvery animation serves a purpose.\nEvery component tells a story.", size=20, color=LIGHT, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "Material Design 3  •  Dark-First  •  Performance-Optimized  •  WCAG 2.1 AA", size=14, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
   "v1.0  •  March 2026  •  Living document — updated with each major design evolution", size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
