#!/usr/bin/env python3
"""Build GigLift IP Strategy PPTX — Trademark, Patent & Risk Analysis."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "data", "docs", "business", "GigLift_IP_Strategy.pptx")
BG = "/Users/bthornley/.gemini/antigravity/brain/d900a025-a940-458f-9e2e-9809b998c050/ip_slide_bg_1772693744829.png"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(150, 150, 170)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
GOLD = RGBColor(234, 179, 8)
DARK_BG = RGBColor(12, 12, 30)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    try:
        slide.shapes.add_picture(BG, Emu(0), Emu(0), W, H)
    except Exception:
        pass


def tx(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(20, 20, 42)
    shape.line.color.rgb = RGBColor(50, 50, 80)
    shape.line.width = Pt(1)
    tx(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
       title, size=16, color=accent, bold=True)
    y = top + Inches(0.55)
    for item in items:
        tx(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3),
           f"• {item}", size=12, color=LIGHT)
        y += Inches(0.28)


def add_table_row(slide, y, cols, sizes, colors=None):
    """Add a row of text items at y position."""
    x = Inches(0.8)
    for i, (text, width) in enumerate(zip(cols, sizes)):
        c = (colors[i] if colors else LIGHT)
        tx(slide, x, y, Inches(width), Inches(0.32), text, size=13, color=c)
        x += Inches(width)


# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(1.3), Inches(11), Inches(0.6),
   "🛡️ INTELLECTUAL PROPERTY STRATEGY", size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(2.1), Inches(11), Inches(1.5),
   "Trademark & Patent\nAnalysis", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
   "Cost Estimates  •  Implementation Plan  •  Risk Assessment", size=22, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "GigLift  •  March 2026", size=16, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 2: Executive Summary
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Executive Summary", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Protecting GigLift's brand and technology", size=16, color=GOLD)
tx(s, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
   "\"GigLift\" is a coined, fanciful mark — the strongest category for trademark protection. "
   "Registration across 2 classes (software + SaaS) is affordable and should be pursued immediately. "
   "The multi-agent AI orchestration system is the strongest patent candidate, though post-Alice "
   "software patent eligibility requires careful claim drafting.", size=17, color=LIGHT)

# Key numbers
for i, (num, label) in enumerate([
    ("$1.4–2.4K", "Trademark\n(2 Classes)"),
    ("$2.1–5.1K", "Provisional Patent"),
    ("5–8 mo", "TM Timeline"),
    ("~$5–8K", "Total Phase 1"),
]):
    x = Inches(0.6) + Inches(3.1) * i
    tx(s, x, Inches(3.8), Inches(2.8), Inches(0.7), num, size=36, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x, Inches(4.5), Inches(2.8), Inches(0.6), label, size=13, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 3: Trademark — Classes
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Trademark Classification", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Two primary classes recommended for filing", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5), "Class 9 — Software", [
    "Downloadable software applications",
    "Mobile apps for event management",
    "AI-powered lead discovery tools",
    "Social media content generators",
    "USPTO fee: $350/class",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.5), "Class 42 — SaaS", [
    "Software-as-a-Service (cloud-based)",
    "AI agent orchestration platform",
    "Content planning & publishing tools",
    "Booking & outreach automation",
    "USPTO fee: $350/class",
], accent=CYAN)

add_card(s, Inches(0.8), Inches(5.6), Inches(5.5), Inches(1.5), "Optional: Class 41 — Education", [
    "Music education services (Instructor mode)",
], accent=MUTED)

add_card(s, Inches(6.8), Inches(5.6), Inches(5.5), Inches(1.5), "Optional: Class 35 — Business Services", [
    "Marketing services for venues/artists",
], accent=MUTED)

# ═══════════════════════════════════════════════
# SLIDE 4: Trademark — Cost Comparison
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Trademark Cost Comparison", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Three filing options — 2 classes (Class 9 + Class 42)", size=16, color=GOLD)

options = [
    ("Option A: Self-Filing", RED, [
        "USPTO base: 2 × $350 = $700",
        "Attorney: $0",
        "Total: $700",
        "",
        "⚠️ ~50% rejection rate",
        "⚠️ Risk of weak claims",
        "⚠️ Office actions costly to fix",
    ]),
    ("Option B: Online Service", AMBER, [
        "Service fee: $500–$1,000",
        "USPTO base: $700",
        "Total: $1,200–$1,700",
        "",
        "✅ Basic professional guidance",
        "✅ Template-based filing",
        "⚠️ Limited strategic advice",
    ]),
    ("Option C: Attorney ★", GREEN, [
        "Comprehensive search: $500–$1,500",
        "Filing + prosecution: $500–$1,500",
        "USPTO base: $700",
        "Office actions: $150–$500 each",
        "Total: $1,850–$4,200",
        "",
        "✅ Strongest protection",
        "✅ Strategic claim drafting",
    ]),
]

for i, (name, color, items) in enumerate(options):
    x = Inches(0.6) + Inches(4.1) * i
    add_card(s, x, Inches(1.8), Inches(3.8), Inches(5.2), name, items, accent=color)

# ═══════════════════════════════════════════════
# SLIDE 5: Trademark — Timeline
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Trademark Timeline", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "5–8 months from filing to registration", size=16, color=GOLD)

timeline = [
    ("Week 1–2", "Clearance Search", "Run TESS search + engage attorney for comprehensive search", AMBER),
    ("Week 2–3", "File Application", "Submit in Classes 9 + 42 with pre-approved descriptions", AMBER),
    ("Month 2–5", "USPTO Examination", "Examiner reviews application, may issue Office Actions", CYAN),
    ("Month 5–6", "Publication", "30-day opposition period in Official Gazette", CYAN),
    ("Month 6–8", "Registration", 'Certificate issues — begin using ® symbol', GREEN),
    ("Year 5–6", "Section 8 Filing", "Declaration of continued use ($225/class)", MUTED),
    ("Year 10", "Renewal", "Section 9 renewal ($300/class)", MUTED),
]

y = Inches(1.8)
for date, title, desc, color in timeline:
    tx(s, Inches(0.8), y, Inches(2), Inches(0.35), date, size=14, color=color, bold=True)
    tx(s, Inches(3), y, Inches(2.8), Inches(0.35), title, size=14, color=WHITE, bold=True)
    tx(s, Inches(5.8), y, Inches(7), Inches(0.35), desc, size=13, color=LIGHT)
    y += Inches(0.58)

# ═══════════════════════════════════════════════
# SLIDE 6: Patent — Eligibility Analysis
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Patent Eligibility Analysis", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Post-Alice test: abstract idea → inventive concept", size=16, color=GOLD)

inventions = [
    ("Multi-Agent AI Orchestration", "⭐⭐⭐⭐", GREEN, [
        "7 agents, autonomous scheduling",
        "Edge-SQLite cross-agent state",
        "Domain-specific prompt chains",
        "Self-healing guardrails",
    ]),
    ("AI Lead Scoring Pipeline", "⭐⭐⭐", AMBER, [
        "Seed→discover→enrich→dedup→score→QC",
        "10+ source web scraping",
        "Novel combination of steps",
        "Risk: individual parts = prior art",
    ]),
    ("Dual-Mode Platform", "⭐", RED, [
        "Performer/Instructor modes",
        "UI/UX feature, not a process",
        "Not patentable",
    ]),
]

for i, (name, stars, color, items) in enumerate(inventions):
    x = Inches(0.6) + Inches(4.1) * i
    add_card(s, x, Inches(1.8), Inches(3.8), Inches(4.5), f"{name}  {stars}", items, accent=color)

tx(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
   "Key: ⭐⭐⭐⭐ = Strong candidate    ⭐⭐⭐ = Moderate    ⭐ = Weak / not patentable", size=12, color=MUTED)

# ═══════════════════════════════════════════════
# SLIDE 7: Patent — Strongest Candidate
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Strongest Patent Candidate", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Multi-Agent AI Orchestration System", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), "What Makes It Patentable", [
    "7 autonomous agents on independent schedules",
    "Edge-SQLite orchestration across serverless infra",
    "Domain-specific AI prompt engineering chains",
    "Self-healing pipeline with guardrails & rate limits",
    "Novel architecture: discovery → score → outreach →",
    "  branding → social → success monitoring",
    "Technical (not business method) framing",
], accent=GREEN)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5), "Risks & Considerations", [
    "Post-Alice: could be seen as 'organizing human",
    "  activity using generic computer components'",
    "Must emphasize technical architecture in claims",
    "Prior art in multi-agent AI systems (growing field)",
    "Patent prosecution takes 2-3 years",
    "Cost: $12K-$24K through grant (micro entity)",
    "Maintenance fees: ~$5K over 20-year lifetime",
], accent=RED)

# ═══════════════════════════════════════════════
# SLIDE 8: Patent — Cost Breakdown
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Patent Cost Breakdown", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "GigLift likely qualifies as micro entity (halved USPTO fees)", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), "Provisional Patent (Rec'd First Step)", [
    "Secures priority date for 12 months",
    "",
    "Attorney drafting: $2,000–$5,000",
    "USPTO filing (micro): $65",
    "Total: $2,065–$5,065",
    "",
    "✅ Low cost, buys time",
    "✅ 12 months to evaluate viability",
], accent=AMBER)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8), "Full Utility Patent (If Warranted)", [
    "20-year protection, enforceable",
    "",
    "Attorney drafting: $8,000–$15,000",
    "USPTO filing + search + exam: ~$825",
    "Office action responses (2-3): $3K–$8K",
    "Issue fee: $500",
    "Total through grant: $12,325–$24,325",
    "Maintenance (lifetime): ~$5,000",
], accent=PURPLE)

# ═══════════════════════════════════════════════
# SLIDE 9: Significant Risks
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Significant IP Risks", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Key risks to consider before filing", size=16, color=RED)

risks = [
    ("⚠️ Trademark Risks", AMBER, [
        "Likelihood of confusion with existing marks containing 'Gig'",
        "Potential opposition from brands in adjacent music/tech spaces",
        "Common-law mark holders may challenge — not all marks in TESS",
        "International protection requires separate filings per country",
        "Failure to use mark within 3 years of registration = cancellation",
    ]),
    ("⚠️ Patent Risks", RED, [
        "Alice rejection: USPTO may classify as abstract business method",
        "Growing prior art in multi-agent AI systems (Google, MSFT, etc.)",
        "Patent prosecution costs can escalate with multiple Office Actions",
        "Competitors may design around claims with minor architectural changes",
        "Enforcement costs ($50K–$500K+) can exceed patent value for startups",
        "Open-source components in stack may complicate claim scope",
    ]),
    ("⚠️ Strategic Risks", MUTED, [
        "IP spend diverts from product development and growth",
        "Patents may not deter well-funded acquirers or competitors",
        "Trade secret protection may be more cost-effective for some IP",
        "International filing costs scale significantly ($3K–$5K/country)",
        "Provisional patent burns 12-month clock — must commit or lose date",
    ]),
]

for i, (name, color, items) in enumerate(risks):
    x = Inches(0.4) + Inches(4.2) * i
    add_card(s, x, Inches(1.8), Inches(3.9), Inches(5.0), name, items, accent=color)

# ═══════════════════════════════════════════════
# SLIDE 10: Budget Allocation
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Recommended Budget", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "From seed round 10% Legal & Contingency ($100K allocation)", size=16, color=GOLD)

# Header
tx(s, Inches(0.8), Inches(1.9), Inches(5), Inches(0.35), "ITEM", size=12, color=MUTED, bold=True)
tx(s, Inches(6.5), Inches(1.9), Inches(2.5), Inches(0.35), "COST RANGE", size=12, color=MUTED, bold=True)
tx(s, Inches(9.5), Inches(1.9), Inches(2.5), Inches(0.35), "PRIORITY", size=12, color=MUTED, bold=True)

budget = [
    ("Trademark — attorney search + filing (2 classes)", "$2,000–$3,000", "🟢 Immediate", WHITE, GREEN),
    ("Provisional patent — multi-agent system", "$3,000–$5,000", "🟡 60 days", WHITE, AMBER),
    ("Domain portfolio (giglift.io, .co, .ai)", "$200", "🟢 Immediate", WHITE, GREEN),
    ("Full utility patent (if provisional succeeds)", "$12,000–$24,000", "🔴 Year 2", WHITE, RED),
    ("International trademark (EU, UK, AU)", "$3,000–$5,000/region", "🔴 Post-Series A", WHITE, RED),
]

y = Inches(2.4)
for item, cost, priority, c1, c2 in budget:
    tx(s, Inches(0.8), y, Inches(5.5), Inches(0.32), item, size=14, color=LIGHT)
    tx(s, Inches(6.5), y, Inches(2.5), Inches(0.32), cost, size=14, color=WHITE, bold=True)
    tx(s, Inches(9.5), y, Inches(3), Inches(0.32), priority, size=14, color=c2)
    y += Inches(0.5)

# Total
tx(s, Inches(0.8), y + Inches(0.2), Inches(5.5), Inches(0.35),
   "TOTAL PHASE 1 SPEND", size=16, color=GOLD, bold=True)
tx(s, Inches(6.5), y + Inches(0.2), Inches(2.5), Inches(0.35),
   "$5,200–$8,200", size=16, color=GOLD, bold=True)
tx(s, Inches(0.8), y + Inches(0.65), Inches(10), Inches(0.3),
   "Remaining legal contingency: $91,800–$94,800", size=13, color=MUTED)

# ═══════════════════════════════════════════════
# SLIDE 11: Action Plan
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Implementation Action Plan", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Prioritized steps for IP protection", size=16, color=GOLD)

steps = [
    ("1", "Run TESS Trademark Search", "Free — search USPTO database for 'GigLift' conflicts in Classes 9/42", GREEN),
    ("2", "Secure Domain Portfolio", "Register giglift.io, .co, .ai to prevent squatting (~$200)", GREEN),
    ("3", "Engage IP Attorney", "Find an attorney specializing in software/SaaS trademarks", AMBER),
    ("4", "File Trademark Application", "2 classes (9 + 42) with pre-approved ID Manual descriptions", AMBER),
    ("5", "Draft Provisional Patent", "Focus on multi-agent orchestration architecture claims", CYAN),
    ("6", "Evaluate Full Utility Patent", "After 6-12 months, assess commercial value and file if warranted", PURPLE),
]

y = Inches(1.8)
for num, title, desc, color in steps:
    tx(s, Inches(0.8), y, Inches(0.5), Inches(0.5), num, size=28, color=color, bold=True)
    tx(s, Inches(1.5), y, Inches(5), Inches(0.35), title, size=18, color=WHITE, bold=True)
    tx(s, Inches(1.5), y + Inches(0.35), Inches(11), Inches(0.3), desc, size=13, color=MUTED)
    y += Inches(0.78)

# ═══════════════════════════════════════════════
# SLIDE 12: Closing
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(2), Inches(11), Inches(1),
   "GigLift", size=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
   "IP Strategy & Protection Plan", size=28, color=GOLD, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
   "Protect the brand. Secure the innovation.", size=20, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
   "This document is for planning purposes only and does not constitute legal advice.\n"
   "Consult a registered IP attorney before filing.", size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
