#!/usr/bin/env python3
"""Build GigLift Coding Standards PPTX — Performance-at-Scale Best Practices."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "data", "docs", "tech", "GigLift_Coding_Standards.pptx")
BG = "/Users/bthornley/.gemini/antigravity/brain/d900a025-a940-458f-9e2e-9809b998c050/coding_standards_bg_1772750926362.png"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(140, 140, 165)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
TEAL = RGBColor(45, 212, 191)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
GOLD = RGBColor(234, 179, 8)
DARK_BG = RGBColor(10, 14, 28)
CARD_BG = RGBColor(16, 20, 38)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    try:
        slide.shapes.add_picture(BG, Emu(0), Emu(0), W, H)
    except Exception:
        pass


def tx(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(35, 45, 70)
    shape.line.width = Pt(1)
    tx(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
       title, size=14, color=accent, bold=True)
    y = top + Inches(0.48)
    for item in items:
        if not item:
            y += Inches(0.08)
            continue
        tx(slide, left + Inches(0.25), y, width - Inches(0.45), Inches(0.26),
           f"• {item}" if not item.startswith(("•", " ", "✅", "❌", "⚠️", "→")) else item, size=11, color=LIGHT)
        y += Inches(0.24)


def header(slide, title, subtitle, subtitle_color=TEAL):
    tx(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.55),
       title, size=32, color=WHITE, bold=True)
    if subtitle:
        tx(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.3),
           subtitle, size=14, color=subtitle_color)


def code_block(slide, left, top, width, height, code, accent=CYAN):
    """Add a styled code block."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(8, 10, 22)
    shape.line.color.rgb = RGBColor(40, 50, 70)
    shape.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(10)
    p.font.color.rgb = accent
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT


# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(0.6), Inches(11), Inches(0.5),
   "📐 ENGINEERING STANDARDS", size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
   "Coding Standards\n& Performance Guidelines", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
   "Best practices for building features that scale", size=20, color=LIGHT, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
   "All new features MUST adhere to these standards before merging to main.", size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
   "GigLift Engineering  •  v1.0  •  March 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Standards Overview", "13 categories — performance at scale + OWASP security as guiding principles")

sections = [
    ("1", "API Route Design", "Request validation, response formats, error handling", CYAN),
    ("2", "Database Queries", "Pagination, indexing, N+1 prevention, connection pooling", CYAN),
    ("3", "React Components", "Server vs. client, code splitting, render optimization", PURPLE),
    ("4", "State Management", "Data fetching patterns, caching, SWR strategies", PURPLE),
    ("5", "AI Agent Guidelines", "Rate limiting, token budgets, prompt efficiency", GREEN),
    ("6", "Security Requirements", "SSRF, input sanitization, auth, secrets management", RED),
    ("7", "OWASP Top 10 Compliance", "Injection, auth, XSS, SSRF, logging — mapped to GigLift controls", RED),
    ("8", "Error Handling", "Boundaries, retries, graceful degradation", AMBER),
    ("9", "TypeScript Conventions", "Strict mode, type exports, naming, file structure", TEAL),
    ("10", "Testing Requirements", "Coverage targets, unit/integration split, mocking", TEAL),
    ("11", "CSS & Styling", "Design tokens, responsive patterns, animation budgets", MUTED),
    ("12", "Performance Budgets", "Bundle size, LCP, CLS, server response time targets", GOLD),
    ("13", "Code Review Checklist", "Pre-merge requirements and CI gates", GOLD),
]

y = Inches(1.4)
for num, title, desc, color in sections:
    tx(s, Inches(0.8), y, Inches(0.5), Inches(0.3), num, size=14, color=color, bold=True)
    tx(s, Inches(1.4), y, Inches(3.5), Inches(0.3), title, size=14, color=WHITE, bold=True)
    tx(s, Inches(5.0), y, Inches(7.5), Inches(0.3), desc, size=12, color=MUTED)
    y += Inches(0.39)

# ═══════════════════════════════════════════════
# SLIDE 3: API Route Design
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "1. API Route Design", "Every route must be safe, validated, and paginated by default")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "REQUEST VALIDATION", [
    "Always validate request body with type guards",
    "Use pickFields() for mass-assignment protection",
    "Validate URL params before database queries",
    "Set Content-Type: application/json explicitly",
    "",
    "✅ DO:",
    "  const { name, email } = pickFields(body, ALLOWED)",
    "",
    "❌ DON'T:",
    "  const data = await req.json() // unvalidated",
    "  await db.execute(sql, [data.anyField])",
], accent=GREEN)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "RESPONSE PATTERNS", [
    "Return consistent { data, error } shape",
    "Always include status codes (200, 201, 400, 404, 500)",
    "Use safeErrorResponse() — never expose internals",
    "Include pagination metadata on list endpoints",
    "",
    "Return shape:",
    "  { data: T[], total, limit, offset, hasMore }",
    "",
    "Error shape:",
    "  { error: 'Human-readable message' }",
    "",
    "Never return stack traces to client",
], accent=CYAN)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "RATE LIMITING & TIMEOUTS", [
    "All external-facing routes: max 60 req/min/user",
    "AI agent routes: max 10 req/min/user",
    "File upload routes: max 5 req/min/user",
    "Set AbortController timeout on all fetches",
    "",
    "Performance targets:",
    "  GET routes: < 200ms p95",
    "  POST routes: < 500ms p95",
    "  AI routes: < 30s timeout",
    "",
    "Use fetchWithRetry() for external calls",
    "Max 3 retries with exponential backoff",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 4: Database Queries
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "2. Database Queries", "Pagination is MANDATORY — unbounded queries are the #1 scaling killer")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5), "PAGINATION (NON-NEGOTIABLE)", [
    "Every list query MUST accept PaginationOptions { limit, offset }",
    "Default page size: 200, max: 1,000 — enforced by clampPageSize()",
    "Always return PaginatedResult<T> with { data, total, limit, offset, hasMore }",
    "Never use SELECT * without LIMIT in production code",
    "Use COUNT(*) in parallel with data query (Promise.all) for total count",
], accent=RED)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5), "QUERY SAFETY", [
    "Always use parameterized queries — NEVER string concatenation",
    "Use escapeLikeQuery() for LIKE/search patterns",
    "Use user_id in every WHERE clause (tenant isolation)",
    "Add ESCAPE '\\\\' to every LIKE query",
    "Validate and sanitize sort fields against an allowlist",
], accent=GREEN)

add_card(s, Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.6), "N+1 PREVENTION", [
    "Use Promise.all() for independent queries, not sequential awaits",
    "Store denormalized data in JSON columns (data field pattern)",
    "Avoid querying inside loops — batch with IN clauses",
    "",
    "✅ const [count, rows] = await Promise.all([countQ, dataQ])",
    "❌ for (const id of ids) { await db.execute(sql, [id]) }",
], accent=CYAN)

add_card(s, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.6), "SCHEMA & INDEXES", [
    "Every table MUST have: id (PK), user_id (indexed), created_at, updated_at",
    "Add indexes on columns used in WHERE, ORDER BY, and JOIN",
    "Use ON CONFLICT for upserts — avoid SELECT-then-INSERT patterns",
    "Keep JSON data columns for flexible schema, index critical fields",
    "",
    "Follow domain module pattern: db/{domain}.ts + db/index.ts barrel",
], accent=PURPLE)

# ═══════════════════════════════════════════════
# SLIDE 5: React Components
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "3. React Components", "Server-first rendering — minimize client JavaScript")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "SERVER vs. CLIENT", [
    "Default to Server Components (no 'use client')",
    "Only add 'use client' when you need:",
    "  → useState, useEffect, event handlers",
    "  → Browser APIs (localStorage, window)",
    "  → Third-party client-only libs",
    "",
    "Server component benefits:",
    "  Zero JS sent to browser",
    "  Direct DB/API access",
    "  Faster initial render",
    "",
    "Split pages: server shell + client islands",
], accent=GREEN)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "COMPONENT STANDARDS", [
    "Max component size: ~300 lines",
    "Extract reusable UI into src/components/",
    "Use TypeScript interfaces for all props",
    "Export named components (not default) from libs",
    "",
    "Naming:",
    "  Components: PascalCase (StatCard.tsx)",
    "  Hooks: camelCase (useAppMode.ts)",
    "  Utils: camelCase (api-client.ts)",
    "",
    "Each component = one file, one concern",
], accent=PURPLE)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "RENDER OPTIMIZATION", [
    "Memoize expensive computations: useMemo()",
    "Stable callbacks with useCallback()",
    "Avoid inline object/array literals in JSX props",
    "Use React.lazy() for heavy components",
    "",
    "❌ Anti-patterns:",
    "  style={{ ...complex }} on every render",
    "  Fetching in useEffect without cleanup",
    "  Missing dependency arrays",
    "  Re-creating arrays/objects on each render",
    "",
    "Prefer CSS classes over inline styles",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 6: State & Data Fetching
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "4. State & Data Fetching", "Minimize waterfalls — fetch data in parallel, cache aggressively")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.6), "DATA FETCHING PATTERNS", [
    "Use fetchWithRetry() from api-client.ts for ALL API calls",
    "Set timeouts on every fetch (default: 15s, AI: 30s)",
    "Fetch on mount with useEffect + AbortController cleanup",
    "Use Promise.all() to parallelize independent fetches on page load",
    "Show skeleton/loading states during fetches — never block the UI",
    "Implement optimistic updates for mutations (update UI before server confirms)",
], accent=CYAN)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.6), "CACHING STRATEGY", [
    "API routes: Set Cache-Control headers for static data",
    "Revalidate on mutation — refetch after POST/PUT/DELETE",
    "Store user-scoped state in component state (not global)",
    "Use localStorage only for UI preferences (mode, theme)",
    "Never cache sensitive data (tokens, PII) client-side",
    "Consider SWR pattern for frequently-updated data",
], accent=GREEN)

add_card(s, Inches(0.6), Inches(4.4), Inches(11.9), Inches(2.4), "WATERFALL PREVENTION (CRITICAL FOR PERFORMANCE)", [
    "✅ Parallel: const [events, leads, stats] = await Promise.all([fetchEvents(), fetchLeads(), fetchStats()])",
    "❌ Sequential: const events = await fetchEvents(); const leads = await fetchLeads(); // 2× slower",
    "",
    "Dashboard example: 5 independent API calls → fetch ALL in parallel → render time = slowest single call",
    "Dependent calls: only waterfall when call B genuinely needs data from call A",
    "",
    "Every page load should resolve in ≤ 2 network round-trips (auth + data)",
], accent=RED)

# ═══════════════════════════════════════════════
# SLIDE 7: AI Agent Guidelines
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "5. AI Agent Guidelines", "Token efficiency, rate limiting, and cost control")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "PROMPT ENGINEERING", [
    "System prompts: < 500 tokens",
    "User prompts: < 2,000 tokens (truncate context)",
    "Response max_tokens: set explicit limit per use case",
    "",
    "Use structured output (JSON mode) when possible",
    "Include few-shot examples inline (1-2 max)",
    "Avoid iterative multi-turn — prefer single-shot",
    "",
    "Cost tracking:",
    "  Log tokens_used per request",
    "  Track cost_per_user per billing period",
    "  Alert at 80% of monthly budget",
], accent=PURPLE)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "RATE LIMITING & QUOTAS", [
    "Enforce user-level monthly quotas (search_quota table)",
    "Use dbIncrementSearchQuota() before AI calls",
    "Return 429 + remaining quota in headers",
    "",
    "Per-user limits:",
    "  Lead scans: plan-based (10/50/unlimited)",
    "  AI generations: 100/mo free, scale w/ plan",
    "  Email sends: 50/mo free tier",
    "",
    "Global guardrails:",
    "  Max concurrent AI calls: 5 per user",
    "  Reject if response time > 30s",
], accent=GREEN)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "AGENT ARCHITECTURE", [
    "Each agent = isolated module in src/lib/agent/",
    "Agents must be stateless (all state in DB)",
    "Use guardrails.ts for input/output validation",
    "Never expose raw AI output to users",
    "",
    "Error handling:",
    "  Wrap all AI calls in try/catch",
    "  Return structured error, not AI hallucination",
    "  Log failures for later analysis",
    "",
    "Testing:",
    "  Mock OpenAI in tests (never hit API)",
    "  Test prompt building separately from API call",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 8: Security Requirements
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "6. Security Requirements", "Every route, every input, every output — verified and sanitized")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "AUTHENTICATION & AUTHORIZATION", [
    "All API routes: verify Clerk auth via currentUser()",
    "Extract userId from auth, never from request body",
    "Admin routes: check role === 'admin' explicitly",
    "Never trust client-side role claims",
    "",
    "Pattern:",
    "  const { userId } = await auth();",
    "  if (!userId) return error(401);",
    "  // userId used in all DB queries",
    "",
    "Never expose other users' data",
    "Always filter by user_id in WHERE clauses",
], accent=RED)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "INPUT SANITIZATION", [
    "Use validateExternalUrl() for any user-supplied URLs",
    "Block SSRF: private IPs, localhost, metadata endpoints",
    "Use escapeLikeQuery() for search inputs",
    "Validate file types + sizes on upload",
    "Use pickFields() to prevent mass assignment",
    "",
    "Never pass raw user input to:",
    "  SQL queries (use parameterized)",
    "  eval() or new Function()",
    "  Shell commands",
    "  AI system prompts (separate user/system)",
], accent=AMBER)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "SECRETS & ENVIRONMENT", [
    "All secrets in .env.local (never committed)",
    "Use process.env only in server-side code",
    "Never prefix secrets with NEXT_PUBLIC_",
    "Rotate API keys on any suspected exposure",
    "",
    "Required env vars (validated at startup):",
    "  TURSO_DATABASE_URL",
    "  TURSO_AUTH_TOKEN",
    "  CLERK_SECRET_KEY",
    "  OPENAI_API_KEY",
    "  STRIPE_SECRET_KEY",
    "  RESEND_API_KEY",
], accent=GREEN)

# ═══════════════════════════════════════════════
# SLIDE 9: OWASP Top 10 Compliance
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "7. OWASP Top 10 Compliance", "All code must address these attack vectors — mapped to GigLift's defenses")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5), "A01: BROKEN ACCESS CONTROL / A02: CRYPTOGRAPHIC FAILURES", [
    "A01: Every API route verifies Clerk auth; user_id in all DB WHERE clauses",
    "A01: Admin routes require explicit role check — never trust client claims",
    "A01: File download routes use an allowlist (ALLOWED_FILES) — no path traversal",
    "A02: All secrets in .env.local, never NEXT_PUBLIC_; HTTPS enforced by Vercel",
    "A02: Stripe webhooks verified via constructEvent() with signing secret",
], accent=RED)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5), "A03: INJECTION / A04: INSECURE DESIGN", [
    "A03: ALL database queries use parameterized statements — zero string concat",
    "A03: LIKE queries sanitized via escapeLikeQuery() with explicit ESCAPE clause",
    "A03: AI prompt injection mitigated: user input isolated from system prompts",
    "A04: Defense-in-depth: Clerk auth → API validation → DB tenant isolation",
    "A04: Security review required for any new external data flow",
], accent=AMBER)

add_card(s, Inches(0.6), Inches(4.2), Inches(3.8), Inches(2.8), "A05–A07", [
    "A05 Security Misconfiguration:",
    "  No debug mode in prod; strict CSP headers",
    "  Environment vars validated at startup",
    "",
    "A06 Vulnerable Components:",
    "  npm audit on every build; pin major versions",
    "  Review before adding new dependencies",
    "",
    "A07 Auth Failures:",
    "  Clerk handles session mgmt + MFA",
    "  No custom password storage ever",
], accent=CYAN)

add_card(s, Inches(4.7), Inches(4.2), Inches(3.8), Inches(2.8), "A08–A09", [
    "A08 Software/Data Integrity:",
    "  Stripe webhooks verified w/ signing secret",
    "  No eval(), no dynamic code execution",
    "  Git-signed commits recommended",
    "",
    "A09 Logging & Monitoring:",
    "  Log auth failures + quota violations",
    "  Never log PII, tokens, or full errors to client",
    "  Vercel analytics + Turso metrics for anomalies",
    "  Structured logging with [Module] prefix",
], accent=GREEN)

add_card(s, Inches(8.8), Inches(4.2), Inches(3.8), Inches(2.8), "A10: SSRF", [
    "validateExternalUrl() on ALL user-supplied URLs",
    "",
    "Blocks:",
    "  Private IPs (10.x, 172.16.x, 192.168.x)",
    "  Localhost (127.0.0.1, ::1, 0.0.0.0)",
    "  Cloud metadata (169.254.169.254)",
    "  Non-HTTP schemes (ftp, file, javascript)",
    "  IPv6 private ranges",
    "  *.internal hostnames",
    "",
    "⚠️ Test coverage: 16 SSRF tests in security.test.ts",
], accent=RED)

# ═══════════════════════════════════════════════
# SLIDE 10: Error Handling
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "8. Error Handling", "Fail gracefully — users should never see raw exceptions")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5), "SERVER-SIDE ERRORS", [
    "Every API route: wrap main logic in try/catch",
    "Use safeErrorResponse() — strips internal details",
    "Log the full error server-side (console.error with context)",
    "Return appropriate HTTP status codes (400 for client, 500 for server)",
    "Never return error.message or error.stack to the client",
], accent=RED)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5), "CLIENT-SIDE ERRORS", [
    "Wrap page-level components in <ErrorBoundary>",
    "Show user-friendly error states, not stack traces",
    "Retry failed fetches with fetchWithRetry() (max 3, exponential backoff)",
    "Implement loading / error / empty states for every data-driven component",
    "Log client errors to console with meaningful context",
], accent=AMBER)

add_card(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.5), "RETRY PATTERNS", [
    "External API calls: 3 retries, delays [500, 1000, 2000]ms",
    "Database calls: single retry on transient errors only",
    "AI calls: 2 retries for rate limits (429), no retry for 400s",
    "Webhook processing: idempotent handlers + dedup by event ID",
    "Abort on timeout — never hang indefinitely",
], accent=CYAN)

add_card(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.5), "GRACEFUL DEGRADATION", [
    "If AI fails → show manual form, not a blank page",
    "If email fails → queue for retry, show warning toast",
    "If external API fails → use cached data with staleness warning",
    "If auth fails → redirect to sign-in, preserve return URL",
    "If DB fails → show maintenance page with support contact",
], accent=GREEN)

# ═══════════════════════════════════════════════
# SLIDE 10: TypeScript Conventions
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "9. TypeScript Conventions", "Strict types, consistent naming, modular file structure")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "TYPE SYSTEM", [
    "strict: true in tsconfig — no exceptions",
    "No 'any' in production code (lint-enforced)",
    "Export interfaces from types.ts for shared models",
    "Use type guards for runtime validation",
    "",
    "✅ Use:",
    "  interface FlyerConfig { ... }",
    "  type AppMode = 'performer' | 'instructor' | ...",
    "  function isLead(x: unknown): x is Lead { ... }",
    "",
    "❌ Avoid:",
    "  as any, @ts-ignore",
    "  Implicit any from missing return types",
], accent=TEAL)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "FILE STRUCTURE", [
    "src/app/ — Next.js pages (route-based)",
    "src/components/ — shared UI components",
    "src/hooks/ — shared React hooks",
    "src/lib/ — business logic & utilities",
    "src/lib/db/ — database modules (by domain)",
    "src/lib/agent/ — AI agent modules",
    "src/__tests__/ — test files",
    "",
    "Module pattern:",
    "  domain.ts — implementation",
    "  index.ts — barrel exports",
    "  domain.test.ts — tests",
], accent=PURPLE)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "NAMING CONVENTIONS", [
    "Files: kebab-case (api-client.ts)",
    "Components: PascalCase (StatCard.tsx)",
    "Functions: camelCase (dbGetAllLeads)",
    "Constants: SCREAMING_SNAKE (MAX_PAGE_SIZE)",
    "Types/Interfaces: PascalCase (LeadFilters)",
    "DB functions prefix: db (dbSaveLead)",
    "",
    "Function prefixes:",
    "  db* — database operations",
    "  fetch* — client API calls",
    "  handle* — event handlers",
    "  validate* — input validation",
], accent=GREEN)

# ═══════════════════════════════════════════════
# SLIDE 11: Testing Requirements
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "10. Testing Requirements", "Tests are the safety net for refactoring and scaling")

add_card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.5), "COVERAGE TARGETS", [
    "Utility functions (security, db helpers): 100% coverage",
    "API routes: integration tests for happy path + error cases",
    "Components: test user interactions, not implementation details",
    "AI agents: mock OpenAI, test prompt building + response parsing",
    "Minimum: every PR must include tests for new logic",
], accent=GREEN)

add_card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5), "TEST STRUCTURE", [
    "Framework: Vitest (vitest.config.mts with @ alias)",
    "Files: *.test.ts alongside the code or in __tests__/",
    "Run: npm test (single run), npm run test:watch (dev mode)",
    "Organization: describe() per function, it() per behavior",
    "Name tests with 'should' pattern: it('should reject private IPs')",
], accent=CYAN)

add_card(s, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.5), "MOCKING GUIDELINES", [
    "Mock external services (OpenAI, Stripe, Resend), never hit real APIs",
    "Mock database with in-memory data structures for unit tests",
    "Use vitest.mock() for module-level mocks",
    "Reset mocks between tests (beforeEach / afterEach)",
    "Test edge cases: empty arrays, null values, timeout errors, rate limits",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.5), "PERFORMANCE TESTS", [
    "Benchmark critical paths before and after changes",
    "Test with realistic data volumes (100+ items, not just 2)",
    "Verify pagination works with large datasets",
    "Ensure no N+1 queries (monitor query count in tests)",
    "Measure rendering time for large lists/grids",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 12: CSS & Styling
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "11. CSS & Styling", "Design tokens, responsive patterns, and animation budgets")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "DESIGN TOKENS", [
    "Use CSS custom properties from globals.css",
    "Never hardcode colors — use var(--color-xxx)",
    "Spacing scale: 4px, 8px, 12px, 16px, 20px, 24px, 32px",
    "Border radius: var(--radius-sm/md/lg)",
    "Elevation: var(--elevation-1/2/3/4)",
    "",
    "M3 Color roles:",
    "  --primary — main actions/brand",
    "  --secondary — supporting elements",
    "  --tertiary — accents",
    "  --error — destructive actions",
    "  --surface — backgrounds",
], accent=PURPLE)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "RESPONSIVE DESIGN", [
    "Mobile-first approach: base styles → media queries up",
    "Breakpoints: 768px (mobile), 1024px (tablet)",
    "Use CSS Grid for 2D layouts, Flexbox for 1D",
    "Min touch target: 44×44px on mobile",
    "Test all pages at 320px, 768px, 1024px, 1440px",
    "",
    "Grid pattern:",
    "  grid-template-columns:",
    "  repeat(auto-fit, minmax(200px, 1fr))",
    "",
    "Stack on mobile, grid on desktop",
], accent=CYAN)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "ANIMATION BUDGETS", [
    "Total animation budget per page: < 16ms/frame",
    "Prefer transform + opacity (GPU-accelerated)",
    "Avoid animating: width, height, top, left, margin",
    "Use will-change sparingly (revokes after animation)",
    "",
    "Motion tokens:",
    "  --motion-duration-fast: 150ms",
    "  --motion-duration-medium: 300ms",
    "  --motion-duration-slow: 500ms",
    "",
    "Prefer CSS transitions over JS animation",
    "Respect prefers-reduced-motion",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 13: Performance Budgets
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "12. Performance Budgets", "Hard limits — new features must not exceed these targets")

# Table header
y = Inches(1.4)
cols = ["Metric", "Target", "Max (Hard Limit)", "How to Measure"]
widths = [3.0, 2.5, 2.5, 3.5]
x = Inches(0.8)
for c, w in zip(cols, widths):
    tx(s, x, y, Inches(w), Inches(0.3), c, size=12, color=TEAL, bold=True)
    x += Inches(w)

budgets = [
    ("Largest Contentful Paint (LCP)", "< 2.0s", "< 2.5s", "Lighthouse / Web Vitals", GREEN),
    ("Cumulative Layout Shift (CLS)", "< 0.05", "< 0.1", "Lighthouse / Web Vitals", GREEN),
    ("First Input Delay (FID)", "< 100ms", "< 200ms", "Lighthouse / Web Vitals", GREEN),
    ("JS Bundle (per route)", "< 150KB gz", "< 250KB gz", "next build output", AMBER),
    ("CSS Bundle (total)", "< 50KB gz", "< 80KB gz", "next build output", AMBER),
    ("API Response (p95)", "< 200ms", "< 500ms", "Vercel analytics", CYAN),
    ("AI Route Response (p95)", "< 10s", "< 30s", "Vercel logs", CYAN),
    ("Database Query", "< 50ms", "< 200ms", "Turso metrics", PURPLE),
    ("Image Size (max per image)", "< 200KB", "< 500KB", "Vercel Blob", PURPLE),
    ("Total Page Weight", "< 1MB", "< 2MB", "Lighthouse", RED),
]

y = Inches(1.85)
for name, target, hard, measure, color in budgets:
    vals = [name, target, hard, measure]
    x = Inches(0.8)
    for v, w in zip(vals, widths):
        tx(s, x, y, Inches(w), Inches(0.28), v, size=12, color=LIGHT if v == name else color)
        x += Inches(w)
    y += Inches(0.38)

tx(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5),
   "⚠️ Any PR that causes a metric to exceed the hard limit MUST include a justification and an optimization plan.",
   size=13, color=AMBER, bold=True)

# ═══════════════════════════════════════════════
# SLIDE 14: Code Review Checklist
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "13. Code Review Checklist", "Every PR must satisfy these requirements before merging to main")

add_card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.5), "MUST PASS (BLOCKING)", [
    "☐ next build completes with zero errors",
    "☐ npm test passes (all tests green)",
    "☐ No TypeScript 'any' or @ts-ignore added",
    "☐ All list queries paginated (PaginatedResult)",
    "☐ Auth check in every API route (userId verified)",
    "☐ User input validated before use",
    "☐ No secrets or PII in code/logs",
    "☐ Error responses use safeErrorResponse()",
    "☐ New functions have tests",
    "☐ Component < 300 lines (or justified)",
], accent=RED)

add_card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.5), "SHOULD VERIFY (ADVISORY)", [
    "☐ Lighthouse score > 90 on affected pages",
    "☐ No N+1 queries introduced",
    "☐ Loading/error/empty states for new UI",
    "☐ Mobile responsive (tested at 320px wide)",
    "☐ CSS uses design tokens, not hardcoded values",
    "☐ fetchWithRetry() for external API calls",
    "☐ Promise.all() for independent parallel calls",
    "☐ AbortController on fetches in useEffect",
    "☐ Semantic HTML + accessibility (alt text, labels)",
    "☐ Console.log removed (use structured logging)",
], accent=AMBER)

add_card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.5), "DOCUMENTATION", [
    "☐ JSDoc on exported functions",
    "☐ README updated for new env vars",
    "☐ Commit messages: conventional format",
    "    feat: / fix: / refactor: / docs: / test:",
    "☐ Breaking changes called out explicitly",
    "☐ Migration script if schema changes",
    "",
    "Commit message examples:",
    "  feat: add social analytics dashboard",
    "  fix: pagination on leads endpoint",
    "  refactor: split db.ts into modules",
    "  test: add security util coverage",
], accent=TEAL)

# ═══════════════════════════════════════════════
# SLIDE 15: Quick Reference
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
header(s, "Quick Reference Card", "Copy-paste patterns for common tasks")

code_block(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(1.5),
    "// API Route Pattern\n"
    "export async function GET(req: NextRequest) {\n"
    "  const { userId } = await auth();\n"
    "  if (!userId) return safeErrorResponse('Unauthorized', 401);\n"
    "  try {\n"
    "    const data = await dbGetItems(userId, { limit: 200 });\n"
    "    return NextResponse.json(data);\n"
    "  } catch { return safeErrorResponse('Failed', 500); }\n"
    "}", CYAN)

code_block(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(1.5),
    "// Database Query Pattern\n"
    "export async function dbGetItems(\n"
    "  userId: string, pagination?: PaginationOptions\n"
    "): Promise<PaginatedResult<Item>> {\n"
    "  const limit = clampPageSize(pagination?.limit);\n"
    "  const offset = pagination?.offset ?? 0;\n"
    "  const [countR, dataR] = await Promise.all([countQ, dataQ]);\n"
    "  return { data, total, limit, offset, hasMore };\n"
    "}", GREEN)

code_block(s, Inches(0.6), Inches(3.2), Inches(5.8), Inches(1.5),
    "// Client Fetch Pattern\n"
    "useEffect(() => {\n"
    "  const ctrl = new AbortController();\n"
    "  Promise.all([\n"
    "    fetchEvents(), fetchLeads(), fetchStats()\n"
    "  ]).then(([events, leads, stats]) => {\n"
    "    setEvents(events); setLeads(leads);\n"
    "  }).catch(e => !ctrl.signal.aborted && setError(e));\n"
    "  return () => ctrl.abort();\n"
    "}, []);", PURPLE)

code_block(s, Inches(6.8), Inches(3.2), Inches(5.8), Inches(1.5),
    "// Component Pattern\n"
    "'use client';\n"
    "interface StatCardProps {\n"
    "  value: string | number;\n"
    "  label: string;\n"
    "  accentColor?: string;\n"
    "}\n"
    "export default function StatCard({\n"
    "  value, label, accentColor = 'var(--accent-purple)'\n"
    "}: StatCardProps) { ... }", AMBER)

code_block(s, Inches(0.6), Inches(5.0), Inches(5.8), Inches(1.2),
    "// Test Pattern\n"
    "describe('validateExternalUrl', () => {\n"
    "  it('should block SSRF to private IPs', () => {\n"
    "    expect(validateExternalUrl('http://10.0.0.1').valid).toBe(false);\n"
    "  });\n"
    "});", RED)

code_block(s, Inches(6.8), Inches(5.0), Inches(5.8), Inches(1.2),
    "// Error Handling Pattern\n"
    "try {\n"
    "  const result = await fetchWithRetry(url, { retries: 3 });\n"
    "  return result.json();\n"
    "} catch (err) {\n"
    "  console.error('[FeatureName] Failed:', err);\n"
    "  return safeErrorResponse('Service unavailable', 503);\n"
    "}", TEAL)

# ═══════════════════════════════════════════════
# SLIDE 16: Closing
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
   "GigLift", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
   "Coding Standards & Performance Guidelines", size=26, color=TEAL, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
   "Build fast. Ship safe. Scale confidently.", size=22, color=LIGHT, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
   "\"Every line of code should be written as if the person maintaining it\nis a violent psychopath who knows where you live.\"", size=14, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.8), Inches(11), Inches(0.3),
   "— Martin Golding (paraphrased)", size=11, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(6.4), Inches(11), Inches(0.5),
   "v1.0  •  March 2026  •  Living document — updated with each major feature release", size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
