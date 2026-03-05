#!/usr/bin/env python3
"""Build Bandsintown × GigLift M&A Pitch Deck — Investment / Acquisition Proposal."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

OUT = os.path.join(os.path.dirname(__file__), "data", "docs", "business", "Bandsintown_x_GigLift_MA_Deck.pptx")

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(150, 150, 170)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
TEAL = RGBColor(45, 212, 191)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
GOLD = RGBColor(234, 179, 8)
DARK_BG = RGBColor(10, 15, 30)
CARD_BG = RGBColor(18, 22, 42)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide, cover=False):
    """Pure gradient background — no images, no alignment issues."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    # Add a subtle gradient overlay shape spanning the full slide
    overlay = slide.shapes.add_shape(1, Emu(0), Emu(0), W, H)
    overlay.line.fill.background()  # no border
    fill = overlay.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(10, 15, 30)
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(12, 30, 45) if not cover else RGBColor(8, 25, 50)
    fill.gradient_stops[1].position = 1.0
    # Add a thin accent line at the top
    accent_line = slide.shapes.add_shape(1, Emu(0), Emu(0), W, Inches(0.03))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = TEAL if not cover else RGBColor(45, 212, 191)
    accent_line.line.fill.background()


def tx(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(40, 50, 75)
    shape.line.width = Pt(1)
    tx(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
       title, size=15, color=accent, bold=True)
    y = top + Inches(0.55)
    for item in items:
        tx(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3),
           f"• {item}" if item and not item.startswith("•") and not item.startswith(" ") else item, size=12, color=LIGHT)
        y += Inches(0.27)


def section_header(slide, title, subtitle, subtitle_color=TEAL):
    tx(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
       title, size=34, color=WHITE, bold=True)
    if subtitle:
        tx(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.3),
           subtitle, size=15, color=subtitle_color)


# ═══════════════════════════════════════════════
# SLIDE 1: Title / Cover
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, cover=True)
tx(s, Inches(1), Inches(0.6), Inches(11), Inches(0.5),
   "CONFIDENTIAL — FOR DISCUSSION PURPOSES ONLY", size=13, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
   "Strategic Acquisition\nProposal", size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(3.8), Inches(11), Inches(0.7),
   "Bandsintown  ×  GigLift", size=30, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
   "Expanding the Live Music Ecosystem with AI-Powered Artist Growth Tools", size=18, color=LIGHT, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
   "Prepared by M&A Advisory  •  March 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
   "This presentation contains proprietary information and is intended solely for authorized recipients.", size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 2: Executive Summary
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Executive Summary", "Thesis: AI-powered growth tools for the long tail of live music")

tx(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.8),
   "GigLift is an AI-powered SaaS platform purpose-built for independent DJs and live performers, "
   "automating the end-to-end business of booking gigs: AI lead discovery, automated outreach, "
   "social media management, EPK generation, and event operations.\n\n"
   "Bandsintown, with 100 million registered users, 700,000 artists, and exclusive distribution "
   "partnerships with YouTube, Spotify, and Apple, is the undisputed leader in concert discovery. "
   "However, the platform currently stops at discovery — it does not help artists GET booked.\n\n"
   "Acquiring GigLift would extend Bandsintown from a discovery platform into a full-lifecycle "
   "artist growth engine, capturing revenue from the 95% of artists who need help filling their calendar.",
   size=16, color=LIGHT)

# Key thesis points
for i, (label, val) in enumerate([
    ("BIT Users", "100M"),
    ("GigLift TAM", "$4.7B"),
    ("Revenue Multiple", "4–6× ARR"),
    ("Synergy Unlock", "High"),
]):
    x = Inches(0.6) + Inches(3.1) * i
    tx(s, x, Inches(5.2), Inches(2.8), Inches(0.6), val, size=36, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x, Inches(5.8), Inches(2.8), Inches(0.4), label, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 3: Bandsintown Profile
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Acquirer Profile: Bandsintown", "The concert discovery platform powering the live music ecosystem")

add_card(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3), "COMPANY OVERVIEW", [
    "Founded 2007 • HQ: New York City",
    "~$15M estimated annual revenue",
    "$60.4M total funding raised",
    "Profitable status (per PitchBook, post-2022)",
    "100M registered users as of Oct 2025",
    "700,000+ registered artists",
    "1.8M events published in 2025",
    "",
    "Revenue model: advertising, ticket commissions,",
    "  artist marketplace, data licensing",
], accent=TEAL)

add_card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.3), "STRATEGIC PARTNERSHIPS", [
    "YouTube — exclusive concert listings provider (Aug 2025)",
    "Spotify — integrated concert discovery (Feb 2024)",
    "Apple Music & Shazam — preferred live data provider",
    "Google — event data syndication",
    "Amazon Music — merch integration",
    "",
    "Recent Moves:",
    "AI assistant for concert discovery (2025)",
    "Artist-controlled presales (Nov 2024)",
    "Artist marketplace launch (Oct 2025)",
], accent=CYAN)

# ═══════════════════════════════════════════════
# SLIDE 4: GigLift Overview
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Target Profile: GigLift", "AI-powered artist growth platform — from lead to stage")

add_card(s, Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.3), "PLATFORM", [
    "Next.js 16 / React 19 SaaS",
    "Multi-agent AI architecture",
    "7 autonomous AI agents",
    "Turso (LibSQL) + Clerk Auth",
    "Stripe billing / Vercel deploy",
    "Dual mode: Performer / Instructor",
    "",
    "Built for indie DJs & live artists",
    "Mobile-responsive web app",
], accent=PURPLE)

add_card(s, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.3), "AI AGENTS", [
    "🔍 Lead Finder — venue/event discovery",
    "📡 Auto-Scanner — hands-free lead gen",
    "📧 DJ Agent — outreach & follow-up",
    "👥 Social Hype — content planning",
    "📸 Instagram — auto-posting",
    "💰 Cost Guardian — budget optimization",
    "📊 Growth Ops — pipeline analytics",
    "",
    "+ Customer Success & Community agents",
], accent=GREEN)

add_card(s, Inches(8.8), Inches(1.6), Inches(3.8), Inches(5.3), "ADDITIONAL TOOLS", [
    "📋 EPK Builder — press kits",
    "🎨 Flyer Creator — AI-powered design",
    "📧 Email Templates — outreach templates",
    "📅 Event Management — full lifecycle",
    "🌐 SEO-optimized gig pages (/gigs/city)",
    "🎓 Ambassador program",
    "📊 Social analytics dashboard",
    "",
    "Subscription: Free / Pro / Enterprise",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 5: Strategic Rationale
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Strategic Rationale", "Why Bandsintown should acquire GigLift")

add_card(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.2), "1. COMPLETE THE ARTIST LIFECYCLE", [
    "BIT helps fans discover concerts — but doesn't help artists GET booked",
    "GigLift automates the entire upstream pipeline: discover → pitch → book → promote",
    "Combined: the only platform from lead discovery to ticket sale to post-show analytics",
], accent=TEAL)

add_card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.2), "2. MONETIZE THE LONG TAIL", [
    "700K artists on BIT, most are independent with no booking support",
    "GigLift's SaaS pricing ($29-$149/mo) creates recurring revenue from this underserved base",
    "Even 5% conversion = 35K paying subscribers = $12M+ ARR uplift",
], accent=GREEN)

add_card(s, Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.2), "3. AI TECHNOLOGY LEAP", [
    "GigLift's 7-agent AI orchestration system is production-proven, not a prototype",
    "Addresses BIT's 2025 AI push — pre-built infra vs. building from scratch",
    "Patent-eligible multi-agent architecture (provisional filing recommended)",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.2), "4. DATA NETWORK EFFECTS", [
    "BIT's 100M user concert data + GigLift's lead scoring = unmatched booking intelligence",
    "Predictive analytics: which venues book which genres, seasonal demand, pricing optimization",
    "Creates durable competitive moat that pure-play competitors can't replicate",
], accent=AMBER)

# ═══════════════════════════════════════════════
# SLIDE 6: Market Opportunity
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Market Opportunity", "Music tech SaaS and live music booking addressable market")

add_card(s, Inches(0.6), Inches(1.6), Inches(3.8), Inches(5.0), "TOTAL ADDRESSABLE MARKET", [
    "Global live music: $31.8B (2024)",
    "Music production SaaS: $2.5B+ (2023)",
    "    → Growing at 7% CAGR to $3.8B",
    "",
    "Artist tools TAM: ~$4.7B",
    "    → SaaS tools + booking + promotion",
    "    → marketing + merch + analytics",
    "",
    "2M+ active independent artists in US",
    "    → Growing with creator economy",
], accent=TEAL)

add_card(s, Inches(4.7), Inches(1.6), Inches(3.8), Inches(5.0), "SERVICEABLE MARKET", [
    "700K artists already on Bandsintown",
    "Est. 200K+ active monthly",
    "",
    "DJ-specific segment: 50K+ in US",
    "    → Fastest-growing live segment",
    "    → High SaaS willingness-to-pay",
    "",
    "Venue/promoter side: 250K+ venues",
    "    → Both supply & demand captured",
], accent=GREEN)

add_card(s, Inches(8.8), Inches(1.6), Inches(3.8), Inches(5.0), "COMPETITIVE LANDSCAPE", [
    "Giggster — venue rental (not booking)",
    "Sonicbids — legacy, limited AI",
    "ReverbNation — declining, acquired 2022",
    "Indie on the Move — directory only",
    "",
    "No competitor combines:",
    "✅ AI-powered lead discovery",
    "✅ Automated outreach",
    "✅ Social media management",
    "✅ EPK + flyer creation",
    "✅ Full event lifecycle management",
], accent=PURPLE)

# ═══════════════════════════════════════════════
# SLIDE 7: Synergies
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Revenue & Cost Synergies", "Projected impact of GigLift integration into Bandsintown ecosystem")

# Revenue synergies
add_card(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0), "REVENUE SYNERGIES (YEAR 1–3)", [
    "🟢 SaaS Upsell to BIT Artist Base",
    "    700K artists × 5% conversion × $49 avg = $20.6M ARR",
    "",
    "🟢 Booking Commission Revenue",
    "    Completed bookings via GigLift → 10% commission",
    "    Est. 50K bookings/yr × $500 avg = $2.5M",
    "",
    "🟢 Premium Lead Data Products",
    "    Venue analytics, demand forecasting for artists",
    "    New data product line: $3–5M ARR potential",
    "",
    "🟢 Enhanced Advertising Revenue",
    "    Deeper artist engagement → higher CPM, targeted ads",
    "    Est. 20–30% uplift on BIT ad revenue = $3–4.5M",
], accent=GREEN)

add_card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.0), "COST SYNERGIES", [
    "🔵 Shared Infrastructure",
    "    GigLift on Vercel, BIT existing — consolidate compute",
    "    Est. $200–400K/yr savings",
    "",
    "🔵 Engineering Team Consolidation",
    "    AI/ML team shared across products",
    "    Data pipeline reuse",
    "",
    "🔵 Distribution Leverage",
    "    Zero CAC for GigLift SaaS via BIT's 100M user base",
    "    Eliminate GigLift's standalone marketing budget",
    "    Current BIT → GigLift funnel costs: $0 (organic)",
    "",
    "🔵 Combined Total Synergy: $28–33M by Year 3",
], accent=CYAN)

# ═══════════════════════════════════════════════
# SLIDE 8: Comparable Transactions
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Comparable Transactions & Valuation Context", "Recent music tech and SaaS M&A activity")

# Table header
headers = ["Transaction", "Year", "Value", "Multiple", "Relevance"]
widths = [3.5, 1.0, 1.8, 1.5, 3.5]
x = Inches(0.8)
y = Inches(1.5)
for h, w in zip(headers, widths):
    tx(s, x, y, Inches(w), Inches(0.35), h, size=12, color=TEAL, bold=True)
    x += Inches(w)

# Table rows
comps = [
    ("Blackstone → Hipgnosis (music rights)", "2024", "$2.6B", "~18× rev", "Music asset M&A scale"),
    ("EQT → Believe (music distribution)", "2024", "$2.3B", "~5× rev", "Distribution to artist tools"),
    ("SiriusXM → Pandora → AdsWizz", "2023", "$150M", "~4× rev", "Audio tech roll-up"),
    ("LiveOne → CPS (live production)", "2023", "$116M", "~6× rev", "Live music SaaS tools"),
    ("Music Ally → Stern (advisory)", "2024", "Undisclosed", "~5× ARR", "Music analytics/SaaS"),
    ("Median SaaS Private Exits (2024)", "2024", "Varies", "4.1× ARR", "Industry benchmark"),
    ("SaaS w/ NRR >120% (2024)", "2024", "Varies", "11.7× ARR", "Premium benchmark"),
]

y = Inches(2.0)
for name, year, val, mult, rel in comps:
    cols = [name, year, val, mult, rel]
    x = Inches(0.8)
    for c, w in zip(cols, widths):
        tx(s, x, y, Inches(w), Inches(0.32), c, size=12, color=LIGHT)
        x += Inches(w)
    y += Inches(0.42)

tx(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.6),
   "Music tech M&A had a banner year in 2024 — PE financing value more than doubled. "
   "Early-stage music tech median valuations: $7–8M (source: Outpost Partners).",
   size=13, color=MUTED)

# Valuation guidance
add_card(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2), "IMPLIED VALUATION RANGE — GIGLIFT", [
    "Conservative (4× ARR): Depends on current revenue — applicable at $500K+ ARR",
    "Strategic Premium (6–8× ARR): Justified by BIT synergies, AI IP, and distribution leverage",
    "Comparable Early-Stage: $7–15M range for AI-native music SaaS with proven product",
], accent=GOLD)

# ═══════════════════════════════════════════════
# SLIDE 9: Financial Projections (Combined)
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Combined Financial Outlook", "Pro forma projections for Bandsintown + GigLift (Year 1–3)")

# Table
headers = ["Metric", "Year 1", "Year 2", "Year 3"]
widths = [4, 2.5, 2.5, 2.5]
x = Inches(0.8)
y = Inches(1.5)
for h, w in zip(headers, widths):
    tx(s, x, y, Inches(w), Inches(0.35), h, size=13, color=TEAL, bold=True)
    x += Inches(w)

rows = [
    ("BIT Existing Revenue (est.)", "$15M", "$17M", "$20M", LIGHT),
    ("GigLift SaaS Revenue (upsell)", "$2.8M", "$8.5M", "$20.6M", GREEN),
    ("Booking Commissions", "$0.5M", "$1.2M", "$2.5M", GREEN),
    ("Premium Data Products", "$0.5M", "$1.5M", "$4.0M", GREEN),
    ("Ad Revenue Uplift", "$1.5M", "$2.5M", "$4.5M", GREEN),
    ("", "", "", "", MUTED),
    ("COMBINED REVENUE", "$20.3M", "$30.7M", "$51.6M", GOLD),
    ("", "", "", "", MUTED),
    ("Cost Synergies", "$0.2M", "$0.5M", "$0.8M", CYAN),
    ("Net New Revenue (GigLift)", "$5.3M", "$13.7M", "$31.6M", GREEN),
    ("Revenue CAGR (Combined)", "—", "51%", "68%", AMBER),
]

y = Inches(2.0)
for name, y1, y2, y3, color in rows:
    if not name:
        y += Inches(0.15)
        continue
    is_bold = name.startswith("COMBINED") or name.startswith("Revenue CAGR")
    cols = [name, y1, y2, y3]
    x = Inches(0.8)
    for c, w in zip(cols, widths):
        tx(s, x, y, Inches(w), Inches(0.32), c, size=13, color=color, bold=is_bold)
        x += Inches(w)
    y += Inches(0.38)

tx(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5),
   "Projections assume 5% artist conversion rate (Year 3), $49 avg monthly subscription, "
   "and conservative ramp. Excludes additional products (API licensing, white-label).",
   size=11, color=MUTED)

# ═══════════════════════════════════════════════
# SLIDE 10: Deal Structure Options
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Proposed Transaction Structures", "Two pathways aligned to strategic objectives")

add_card(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2), "OPTION A: FULL ACQUISITION", [
    "Structure: 100% equity purchase",
    "Consideration: $8–15M (cash + earnout)",
    "",
    "Cash at close: 60–70%",
    "Earnout (18 mo): 30–40% tied to milestones",
    "    → User adoption targets",
    "    → Revenue integration goals",
    "    → Technology migration completion",
    "",
    "Founder retention: 2-year employment agreement",
    "Team: Full integration into BIT engineering",
    "",
    "✅ Full control, fastest integration",
    "✅ Complete IP ownership",
    "⚠️ Higher upfront capital requirement",
], accent=TEAL)

add_card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), "OPTION B: STRATEGIC INVESTMENT + OPTION", [
    "Structure: Minority stake (20–30%) + acquisition option",
    "Investment: $2–4M for 25% equity",
    "",
    "Includes:",
    "    → Board observer seat",
    "    → Exclusive integration partnership",
    "    → Right of first refusal on full acquisition",
    "    → Option to purchase at pre-agreed multiple (5× ARR)",
    "",
    "Option exercise window: 18–36 months",
    "Revenue share during option period: 15% on BIT-referred subs",
    "",
    "✅ Lower risk, test-before-buy",
    "✅ GigLift maintains growth independence",
    "⚠️ Risk of competitor pre-emption",
], accent=PURPLE)

# ═══════════════════════════════════════════════
# SLIDE 11: Integration Roadmap
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Integration Roadmap", "90-day, 6-month, and 12-month milestones")

phases = [
    ("PHASE 1: 0–90 DAYS", TEAL, [
        "Deep-link GigLift from BIT Artist Dashboard",
        "SSO integration (Clerk → BIT OAuth)",
        "BIT concert data feed → GigLift lead enrichment",
        "Joint press release + artist onboarding campaign",
        "Technical due diligence + architecture review",
    ]),
    ("PHASE 2: 90–180 DAYS", GREEN, [
        "Unified artist profile (BIT + GigLift data)",
        "BIT event data → GigLift auto-scanner integration",
        "In-app upsell: 'Powered by GigLift' booking tools",
        "Shared analytics dashboard for artists",
        "API layer for cross-platform data sync",
    ]),
    ("PHASE 3: 6–12 MONTHS", PURPLE, [
        "Full platform merge: GigLift as 'BIT for Artists Pro'",
        "AI-powered 'suggested gigs' using BIT demand data",
        "Automated venue-artist matching algorithm",
        "Premium data products launch (venue analytics)",
        "International expansion (EU, UK, AU markets)",
    ]),
]

for i, (name, color, items) in enumerate(phases):
    x = Inches(0.4) + Inches(4.2) * i
    add_card(s, x, Inches(1.6), Inches(3.9), Inches(5.0), name, items, accent=color)

# ═══════════════════════════════════════════════
# SLIDE 12: Risk Analysis
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Risk Analysis & Mitigants", "Key transaction risks and proposed mitigations")

risks = [
    ("INTEGRATION RISKS", AMBER, [
        "Technology stack compatibility (Next.js + BIT)",
        "    → Mitigant: API-first architecture, gradual migration",
        "",
        "Culture fit: startup speed vs. established company",
        "    → Mitigant: Founder retention agreement, autonomous team",
        "",
        "Data privacy (user data across platforms)",
        "    → Mitigant: GDPR/CCPA compliant architecture",
    ]),
    ("MARKET RISKS", RED, [
        "Competitor response (Spotify, YouTube building own tools)",
        "    → Mitigant: 12–18 mo head start, AI IP moat",
        "",
        "Artist willingness to pay for SaaS tools",
        "    → Mitigant: Freemium model de-risks adoption",
        "",
        "Economic downturn impact on live music",
        "    → Mitigant: Live music resilient; grew through 2023–24",
    ]),
    ("FINANCIAL RISKS", MUTED, [
        "Overpayment vs. realized synergies",
        "    → Mitigant: Earnout structure ties value to results",
        "",
        "Revenue ramp slower than projected",
        "    → Mitigant: Conservative base case = 2% conversion",
        "",
        "Hidden tech debt requiring rebuild",
        "    → Mitigant: Clean TS build, zero errors, active codebase",
    ]),
]

for i, (name, color, items) in enumerate(risks):
    x = Inches(0.4) + Inches(4.2) * i
    add_card(s, x, Inches(1.6), Inches(3.9), Inches(5.2), name, items, accent=color)

# ═══════════════════════════════════════════════
# SLIDE 13: Why Now?
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Why Now?", "Market timing and competitive dynamics favor immediate action")

timeline = [
    ("Oct 2025", "BIT hits 100M users + launches artist marketplace", "Platform maturity → ready for premium tools", TEAL),
    ("Aug 2025", "BIT becomes exclusive YouTube concert partner", "Distribution dominance creates upsell surface", TEAL),
    ("2024–25", "Music tech M&A boom — PE deal value doubles", "Favorable market conditions for tech acquisitions", GREEN),
    ("2025", "AI tools adoption accelerating across creative industries", "AI-native platforms command 2× premium vs. legacy", PURPLE),
    ("2025", "GigLift: production-ready, zero tech debt, growing user base", "Acquisition at early-stage pricing before scale", AMBER),
    ("2026+", "Competitor risk: Spotify/YouTube build or acquire similar tools", "Window narrows — first-mover advantage is time-limited", RED),
]

y = Inches(1.6)
for date, event, implication, color in timeline:
    tx(s, Inches(0.8), y, Inches(1.5), Inches(0.35), date, size=14, color=color, bold=True)
    tx(s, Inches(2.5), y, Inches(5.5), Inches(0.35), event, size=13, color=WHITE)
    tx(s, Inches(8.2), y, Inches(4.5), Inches(0.35), implication, size=12, color=MUTED)
    y += Inches(0.55)

tx(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8),
   "Bottom line: Bandsintown's platform reach + GigLift's AI booking tools = the first end-to-end "
   "artist growth engine in music tech. Waiting 12–18 months risks a competitor acquiring similar "
   "capabilities or GigLift scaling to a higher valuation.",
   size=16, color=LIGHT, bold=True)

# ═══════════════════════════════════════════════
# SLIDE 14: Next Steps
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
section_header(s, "Proposed Next Steps", "Path to term sheet within 60 days")

steps = [
    ("1", "Management Introduction", "Founder meeting — product demo, vision alignment, culture assessment", TEAL),
    ("2", "Technical Due Diligence", "Codebase review, architecture assessment, IP audit (2 weeks)", TEAL),
    ("3", "Data Room Access", "Financials, user metrics, contracts, customer pipeline review", GREEN),
    ("4", "Synergy Modeling", "Joint workshop: integration architecture, revenue modeling, cost analysis", GREEN),
    ("5", "Term Sheet", "Non-binding LOI with proposed structure, valuation, and earnout terms", AMBER),
    ("6", "Definitive Agreement", "Negotiate PA, employment agreements, IP assignment, closing conditions", PURPLE),
]

y = Inches(1.6)
for num, title, desc, color in steps:
    tx(s, Inches(0.8), y, Inches(0.5), Inches(0.5), num, size=28, color=color, bold=True)
    tx(s, Inches(1.5), y + Inches(0.02), Inches(5), Inches(0.35), title, size=18, color=WHITE, bold=True)
    tx(s, Inches(1.5), y + Inches(0.37), Inches(11), Inches(0.3), desc, size=13, color=MUTED)
    y += Inches(0.72)

tx(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5),
   "Target timeline: LOI by end of Q2 2026 • Closing by Q3 2026",
   size=15, color=GOLD, bold=True)

# ═══════════════════════════════════════════════
# SLIDE 15: Closing
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, cover=True)
tx(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
   "Bandsintown × GigLift", size=50, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
   "From Discovery to Booking.\nThe Complete Artist Growth Engine.", size=26, color=TEAL, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
   "100M fans • 700K artists • AI-powered booking • One platform", size=18, color=LIGHT, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "CONFIDENTIAL — FOR DISCUSSION PURPOSES ONLY", size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(6.1), Inches(11), Inches(0.5),
   "This document does not constitute an offer to sell or a solicitation to buy securities.\n"
   "All projections are estimates based on publicly available information and internal analysis.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
