#!/usr/bin/env python3
"""Build GigLift Year 1 Projected Financial Plan PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "data", "docs", "business", "GigLift_Year1_Financial_Plan.pptx")
BG = "/Users/bthornley/.gemini/antigravity/brain/d900a025-a940-458f-9e2e-9809b998c050/ip_slide_bg_1772693744829.png"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(150, 150, 170)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
GOLD = RGBColor(234, 179, 8)
DARK_BG = RGBColor(12, 12, 30)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    try:
        slide.shapes.add_picture(BG, Emu(0), Emu(0), W, H)
    except Exception:
        pass


def tx(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(20, 20, 42)
    shape.line.color.rgb = RGBColor(50, 50, 80)
    shape.line.width = Pt(1)
    tx(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
       title, size=16, color=accent, bold=True)
    y = top + Inches(0.55)
    for item in items:
        tx(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3),
           f"• {item}" if item else "", size=12, color=LIGHT)
        y += Inches(0.28)


def table_row(slide, y, cols, sizes, colors=None, bold_flags=None):
    """Helper to render a table row."""
    x = Inches(0.8)
    for i, (text, width) in enumerate(zip(cols, sizes)):
        c = (colors[i] if colors else LIGHT)
        b = (bold_flags[i] if bold_flags else False)
        tx(slide, x, y, Inches(width), Inches(0.3), text, size=13, color=c, bold=b)
        x += Inches(width)


# ══════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(1.3), Inches(11), Inches(0.6),
   "📊 YEAR 1 PROJECTED FINANCIAL PLAN", size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(2.1), Inches(11), Inches(1.5),
   "Revenue Projections, Costs\n& Risk Analysis", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
   "Month-by-Month Forecast  •  Q1–Q4 Breakdown  •  Break-Even Analysis", size=22, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "GigLift  •  March 2026", size=16, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# SLIDE 2: Revenue Model Overview
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Revenue Model", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Four subscription tiers — freemium model with strong upgrade incentives", size=16, color=GOLD)

plans = [
    ("🆓 Free", LIGHT, "$0/mo", [
        "5 scans/month", "25 lead storage", "1 region",
        "Basic EPK page", "3 AI flyer backgrounds",
    ]),
    ("⚡ Pro", CYAN, "$33/mo", [
        "50 auto-scans", "Unlimited leads", "2 regions",
        "AI outreach (3 variants)", "Email send & tracking",
        "20 AI backgrounds", "CSV export",
    ]),
    ("🚀 Unlimited", PURPLE, "$79/mo", [
        "Unlimited everything", "Priority lead scoring",
        "Full flyer suite", "Full Social Suite",
        "Analytics dashboard", "Priority support",
    ]),
    ("🏢 Agency", GOLD, "$149/mo", [
        "5 artist profiles", "Everything in Unlimited",
        "Team dashboard", "Bulk outreach",
        "White-label EPK & flyers", "Dedicated account mgr",
    ]),
]

for i, (name, color, price, features) in enumerate(plans):
    x = Inches(0.4) + Inches(3.15) * i
    add_card(s, x, Inches(1.8), Inches(3.0), Inches(3.5), name, features, accent=color)
    tx(s, x, Inches(5.4), Inches(3.0), Inches(0.4), price, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)

tx(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8),
   "💡 Blended average revenue per paying user: ~$52/mo based on projected 50% Pro, 35% Unlimited, 15% Agency mix.",
   size=13, color=MUTED)


# ══════════════════════════════════════════════
# SLIDE 3: Year 1 Subscriber Forecast
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Year 1 Subscriber Forecast", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Conservative growth model — free-to-paid conversion rate: 8–12%", size=16, color=GOLD)

# Q-by-Q table
cols_h = ["QUARTER", "TOTAL USERS", "FREE", "PAID SUBS", "MRR", "CONVERSION"]
sizes = [1.8, 2.0, 1.8, 1.8, 2.0, 2.0]
y = Inches(1.9)
table_row(s, y, cols_h, sizes, colors=[MUTED]*6, bold_flags=[True]*6)

rows = [
    ("Q1 (Mar–May)", "500", "460", "40", "$2,080", "8%", CYAN),
    ("Q2 (Jun–Aug)", "1,500", "1,350", "150", "$7,800", "10%", CYAN),
    ("Q3 (Sep–Nov)", "3,500", "3,080", "420", "$21,840", "12%", GREEN),
    ("Q4 (Dec–Feb)", "6,000", "5,280", "720", "$37,440", "12%", GREEN),
]

y = Inches(2.4)
for q, total, free, paid, mrr, conv, color in rows:
    table_row(s, y, [q, total, free, paid, mrr, conv], sizes,
              colors=[WHITE, LIGHT, MUTED, color, GREEN, AMBER],
              bold_flags=[True, False, False, True, True, False])
    y += Inches(0.45)

tx(s, Inches(0.8), y + Inches(0.3), Inches(5), Inches(0.35),
   "YEAR 1 TOTAL REVENUE", size=18, color=GOLD, bold=True)
tx(s, Inches(6.8), y + Inches(0.3), Inches(3), Inches(0.35),
   "~$278,000 ARR", size=18, color=GOLD, bold=True)

tx(s, Inches(0.8), y + Inches(0.9), Inches(11.5), Inches(0.6),
   "Assumptions: Blended ARPU $52/mo (50% Pro/$33, 35% Unlimited/$79, 15% Agency/$149). "
   "Monthly churn rate: 5%. Growth powered by ambassador program, SEO, and DJ community word-of-mouth.",
   size=12, color=MUTED)


# ══════════════════════════════════════════════
# SLIDE 4: Month-by-Month Revenue (Q1–Q2)
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Month-by-Month Revenue — Q1 & Q2", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "First 6 months: building the base with aggressive free-tier growth", size=16, color=GOLD)

cols_h = ["MONTH", "NEW USERS", "TOTAL USERS", "PAID SUBS", "MRR", "CUM. REV"]
sizes = [1.8, 1.6, 1.8, 1.6, 2.0, 2.4]
y = Inches(1.9)
table_row(s, y, cols_h, sizes, colors=[MUTED]*6, bold_flags=[True]*6)

months_q1q2 = [
    ("Mar", "50", "50", "4", "$208", "$208"),
    ("Apr", "120", "170", "14", "$728", "$936"),
    ("May", "180", "350", "28", "$1,456", "$2,392"),
    ("Jun", "250", "600", "54", "$2,808", "$5,200"),
    ("Jul", "400", "1,000", "100", "$5,200", "$10,400"),
    ("Aug", "500", "1,500", "150", "$7,800", "$18,200"),
]

y = Inches(2.4)
for m, new, tot, paid, mrr, cum in months_q1q2:
    c = CYAN if int(paid) < 100 else GREEN
    table_row(s, y, [m, new, tot, paid, mrr, cum], sizes,
              colors=[WHITE, LIGHT, LIGHT, c, GREEN, GOLD],
              bold_flags=[True, False, False, True, True, True])
    y += Inches(0.42)

tx(s, Inches(0.8), y + Inches(0.3), Inches(11.5), Inches(0.6),
   "📈 Key milestone: Reach 100 paid subscribers by end of July (Month 5). "
   "DJ festival season (Jun–Aug) expected to accelerate signups.",
   size=13, color=LIGHT)


# ══════════════════════════════════════════════
# SLIDE 5: Month-by-Month Revenue (Q3–Q4)
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Month-by-Month Revenue — Q3 & Q4", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Months 7–12: scaling toward profitability with ambassador-driven growth", size=16, color=GOLD)

y = Inches(1.9)
table_row(s, y, cols_h, sizes, colors=[MUTED]*6, bold_flags=[True]*6)

months_q3q4 = [
    ("Sep", "600", "2,100", "252", "$13,104", "$31,304"),
    ("Oct", "650", "2,750", "330", "$17,160", "$48,464"),
    ("Nov", "750", "3,500", "420", "$21,840", "$70,304"),
    ("Dec", "700", "4,200", "504", "$26,208", "$96,512"),
    ("Jan '27", "800", "5,000", "600", "$31,200", "$127,712"),
    ("Feb '27", "1,000", "6,000", "720", "$37,440", "$165,152"),
]

y = Inches(2.4)
for m, new, tot, paid, mrr, cum in months_q3q4:
    table_row(s, y, [m, new, tot, paid, mrr, cum], sizes,
              colors=[WHITE, LIGHT, LIGHT, GREEN, GREEN, GOLD],
              bold_flags=[True, False, False, True, True, True])
    y += Inches(0.42)

tx(s, Inches(0.8), y + Inches(0.3), Inches(11.5), Inches(0.6),
   "🎯 Year 1 total cumulative revenue: ~$165K. MRR exits Year 1 at ~$37K/mo ($449K ARR run-rate).",
   size=13, color=LIGHT)


# ══════════════════════════════════════════════
# SLIDE 6: Infrastructure Costs by Quarter
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Infrastructure Costs — Year 1", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "All third-party service costs by quarter — includes hosting, AI, payments, email, analytics", size=16, color=GOLD)

cols_h = ["SERVICE", "Q1", "Q2", "Q3", "Q4", "YEAR 1 TOTAL"]
sizes = [2.6, 1.6, 1.6, 1.6, 1.6, 2.4]
y = Inches(1.9)
table_row(s, y, cols_h, sizes, colors=[MUTED]*6, bold_flags=[True]*6)

infra = [
    ("Vercel (Pro)", "$60", "$60", "$60", "$75", "$255", AMBER),
    ("Clerk (Auth)", "$0", "$0", "$0", "$60", "$60", GREEN),
    ("Turso (Database)", "$0", "$87", "$87", "$87", "$261", AMBER),
    ("OpenAI API (7 agents)", "$30", "$90", "$200", "$350", "$670", RED),
    ("Resend (Email)", "$0", "$60", "$60", "$60", "$180", AMBER),
    ("Stripe Fees (2.9%+30¢)", "$17", "$143", "$455", "$823", "$1,438", AMBER),
    ("Plausible (Analytics)", "$27", "$27", "$57", "$57", "$168", GREEN),
    ("Domains", "$9", "$9", "$9", "$9", "$36", GREEN),
    ("Vercel Blob (Storage)", "$0", "$0", "$5", "$10", "$15", GREEN),
]

y = Inches(2.4)
for name, q1, q2, q3, q4, total, color in infra:
    table_row(s, y, [name, q1, q2, q3, q4, total], sizes,
              colors=[LIGHT, MUTED, MUTED, MUTED, MUTED, color],
              bold_flags=[False, False, False, False, False, True])
    y += Inches(0.38)

y += Inches(0.1)
table_row(s, y, ["TOTAL INFRA COSTS", "$143", "$476", "$933", "$1,531", "$3,083"], sizes,
          colors=[GOLD, GOLD, GOLD, GOLD, GOLD, GOLD],
          bold_flags=[True, True, True, True, True, True])


# ══════════════════════════════════════════════
# SLIDE 7: One-Time & Operating Costs
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "One-Time & Operating Costs", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Non-recurring launch costs and ongoing operational expenses", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), "💰 One-Time Launch Costs", [
    "Trademark registration (2 classes): $2,500",
    "Provisional patent filing: $3,000",
    "Legal (Terms of Service, Privacy): $1,500",
    "Domain portfolio (extra TLDs): $100",
    "Apple Developer Program: $99",
    "Google Play Console: $25",
    "",
    "Brand / design refresh: $500",
    "Penetration testing (if needed Y1): $2,000",
    "",
    "TOTAL ONE-TIME: ~$9,724",
], accent=AMBER)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8), "📋 Ongoing Operating Costs", [
    "Founder salary / draw: $0 (reinvested Y1)",
    "Contractor / freelance support: $0–3,000",
    "Accounting / bookkeeping: $150/mo",
    "Business insurance (E&O): $100/mo",
    "Marketing / ads budget (Q3+): $500–2,000/mo",
    "Community tools (Discord Nitro): $10/mo",
    "",
    "TOTAL OPERATING (YEAR 1): ~$6,000–$24,000",
    "",
    "NOTE: No full-time hires planned in Year 1",
], accent=CYAN)


# ══════════════════════════════════════════════
# SLIDE 8: Profit & Loss Summary
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Projected Profit & Loss — Year 1", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Conservative scenario — $0 founder salary, reinvesting all profit", size=16, color=GOLD)

cols_h = ["", "Q1", "Q2", "Q3", "Q4", "YEAR 1"]
sizes = [3.0, 1.6, 1.8, 1.8, 1.8, 2.0]
y = Inches(1.9)
table_row(s, y, cols_h, sizes, colors=[MUTED]*6, bold_flags=[True]*6)

pnl = [
    ("REVENUE", "$2,392", "$15,808", "$52,104", "$94,848", "$165,152", GREEN, True),
    ("", "", "", "", "", "", WHITE, False),
    ("Infrastructure Costs", "($143)", "($476)", "($933)", "($1,531)", "($3,083)", AMBER, False),
    ("Stripe Processing", "($17)", "($143)", "($455)", "($823)", "($1,438)", AMBER, False),
    ("One-Time Costs", "($9,724)", "$0", "$0", "$0", "($9,724)", RED, False),
    ("Operating Costs", "($1,300)", "($2,000)", "($3,500)", "($5,200)", "($12,000)", AMBER, False),
    ("", "", "", "", "", "", WHITE, False),
    ("TOTAL COSTS", "($11,184)", "($2,619)", "($4,888)", "($7,554)", "($26,245)", RED, True),
    ("", "", "", "", "", "", WHITE, False),
    ("NET INCOME", "($8,792)", "$13,189", "$47,216", "$87,294", "$138,907", GREEN, True),
    ("NET MARGIN", "—", "83%", "91%", "92%", "84%", CYAN, True),
]

y = Inches(2.3)
for name, q1, q2, q3, q4, yr, color, bold in pnl:
    if name == "":
        y += Inches(0.12)
        continue
    table_row(s, y, [name, q1, q2, q3, q4, yr], sizes,
              colors=[WHITE if bold else LIGHT, color, color, color, color, color],
              bold_flags=[bold]*6)
    y += Inches(0.38)


# ══════════════════════════════════════════════
# SLIDE 9: Break-Even Analysis
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Break-Even Analysis", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "When does GigLift become self-sustaining?", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.8), "📉 Fixed Monthly Costs (at launch)", [
    "Vercel Pro: $20",
    "Turso Scaler: $29",
    "Resend Pro: $20",
    "Plausible: $9–19",
    "Accounting: $150",
    "Insurance: $100",
    "Domains: ~$3",
    "",
    "TOTAL FIXED: ~$331–341/mo",
], accent=AMBER)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.8), "🎯 Break-Even Point", [
    "Blended ARPU: $52/mo",
    "Variable cost per user (AI + Stripe): ~$4.50",
    "Contribution margin per sub: ~$47.50",
    "",
    "Break-even subs: 341 ÷ 47.50 = ~8 paid subs",
    "",
    "EXPECTED BREAK-EVEN: Month 1 (March '26)",
    "",
    "Note: Excludes one-time launch costs",
    "Full payback (incl. one-time): Month 4",
], accent=GREEN)

tx(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.0),
   "💡 GigLift's cost structure is extremely lean. With near-zero marginal costs from the serverless "
   "architecture, break-even is achieved with just 8 paying subscribers. The one-time launch costs "
   "($9.7K) are fully paid back by Month 4 from cumulative profits.",
   size=13, color=LIGHT)


# ══════════════════════════════════════════════
# SLIDE 10: Cash Flow Projection
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Cash Flow Projection", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Monthly cash position assuming $0 starting capital", size=16, color=GOLD)

cols_h = ["MONTH", "CASH IN", "CASH OUT", "NET", "CASH BALANCE"]
sizes = [2.2, 2.2, 2.2, 2.2, 2.6]
y = Inches(1.9)
table_row(s, y, cols_h, sizes, colors=[MUTED]*5, bold_flags=[True]*5)

cashflow = [
    ("Mar", "$208", "($10,000)", "($9,792)", "($9,792)", RED),
    ("Apr", "$728", "($500)", "$228", "($9,564)", RED),
    ("May", "$1,456", "($550)", "$906", "($8,658)", RED),
    ("Jun", "$2,808", "($600)", "$2,208", "($6,450)", AMBER),
    ("Jul", "$5,200", "($700)", "$4,500", "($1,950)", AMBER),
    ("Aug", "$7,800", "($800)", "$7,000", "$5,050", GREEN),
    ("Sep", "$13,104", "($1,100)", "$12,004", "$17,054", GREEN),
    ("Oct", "$17,160", "($1,300)", "$15,860", "$32,914", GREEN),
    ("Nov", "$21,840", "($1,500)", "$20,340", "$53,254", GREEN),
    ("Dec", "$26,208", "($1,700)", "$24,508", "$77,762", GREEN),
]

y = Inches(2.4)
for m, cin, cout, net, bal, color in cashflow:
    table_row(s, y, [m, cin, cout, net, bal], sizes,
              colors=[WHITE, GREEN, RED, color, color],
              bold_flags=[True, False, False, True, True])
    y += Inches(0.38)

tx(s, Inches(0.8), y + Inches(0.15), Inches(11.5), Inches(0.6),
   "📌 Cash-positive by Month 6 (August). Maximum cash deficit: ~$9.8K in Month 1. Self-funded with <$10K startup capital.",
   size=13, color=LIGHT)


# ══════════════════════════════════════════════
# SLIDE 11: Revenue Scenarios
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Revenue Scenarios — Year 1 Exit Rate", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Best case, base case, and worst case projections", size=16, color=GOLD)

scenarios = [
    ("🥇 Best Case", GREEN, [
        "8,000 total users by Month 12",
        "15% conversion rate",
        "1,200 paid subscribers",
        "MRR: $62,400",
        "ARR: ~$749K",
        "Year 1 revenue: ~$260K",
        "",
        "Driven by: viral ambassador program,",
        "strategic partnerships, festival season",
    ]),
    ("📊 Base Case", CYAN, [
        "6,000 total users by Month 12",
        "12% conversion rate",
        "720 paid subscribers",
        "MRR: $37,440",
        "ARR: ~$449K",
        "Year 1 revenue: ~$165K",
        "",
        "Driven by: organic growth, SEO,",
        "word-of-mouth, ambassador program",
    ]),
    ("⚠️ Worst Case", RED, [
        "3,000 total users by Month 12",
        "6% conversion rate",
        "180 paid subscribers",
        "MRR: $9,360",
        "ARR: ~$112K",
        "Year 1 revenue: ~$52K",
        "",
        "Still cash-positive after Month 8",
        "Pivot/adjust pricing if needed",
    ]),
]

for i, (name, color, items) in enumerate(scenarios):
    x = Inches(0.6) + Inches(4.1) * i
    add_card(s, x, Inches(1.8), Inches(3.8), Inches(5.0), name, items, accent=color)


# ══════════════════════════════════════════════
# SLIDE 12: Key Risks
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Key Financial Risks", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Risk factors ranked by impact and likelihood — with mitigations", size=16, color=GOLD)

cols_h = ["RISK", "IMPACT", "LIKELIHOOD", "MITIGATION"]
sizes = [3.2, 1.4, 1.5, 5.3]
y = Inches(1.9)
table_row(s, y, cols_h, sizes, colors=[MUTED]*4, bold_flags=[True]*4)

risks = [
    ("Slow user acquisition", "HIGH", "Medium", "Ambassador program, DJ community partnerships, content marketing", RED),
    ("High churn rate (>8%/mo)", "HIGH", "Medium", "Improve onboarding, engagement loops, feature value prop", RED),
    ("OpenAI price increase", "Medium", "Low", "Switch to open-source models (Llama/Mistral) for 80% savings", AMBER),
    ("Turso scale jump ($29→$499)", "Medium", "Low", "Monitor usage, optimize queries, consider PlanetScale", AMBER),
    ("Market saturation / competitors", "HIGH", "Medium", "First-mover advantage, niche focus on DJs specifically", RED),
    ("Stripe fee increases", "Low", "Low", "Minor impact — pass through to pricing if needed", GREEN),
    ("Vercel usage overage", "Low", "Medium", "ISR/static generation, CDN caching, optimize API calls", GREEN),
    ("Regulatory / compliance costs", "Medium", "Low", "Budget $5K reserve for SOC 2 prep if needed in Y2", AMBER),
]

y = Inches(2.4)
for name, impact, likelihood, mitigation, color in risks:
    ic = RED if impact == "HIGH" else (AMBER if impact == "Medium" else GREEN)
    lc = AMBER if likelihood == "Medium" else GREEN
    table_row(s, y, [name, impact, likelihood, mitigation], sizes,
              colors=[LIGHT, ic, lc, MUTED],
              bold_flags=[False, True, False, False])
    y += Inches(0.40)


# ══════════════════════════════════════════════
# SLIDE 13: Key Metrics & KPIs
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Key Metrics to Track", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Financial KPIs and targets for Year 1 — tracked monthly", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(3.8), Inches(4.5), "💰 Revenue Metrics", [
    "MRR growth rate: >20% MoM",
    "ARPU: >$50/mo",
    "LTV: >$500 (10+ month retention)",
    "LTV/CAC ratio: >3:1",
    "Net revenue retention: >100%",
    "",
    "TARGET: $37K MRR by Month 12",
], accent=GREEN)

add_card(s, Inches(4.9), Inches(1.8), Inches(3.8), Inches(4.5), "📈 Growth Metrics", [
    "Free-to-paid conversion: >10%",
    "Monthly churn rate: <5%",
    "New signups/month: >500 by Q4",
    "Activation rate (1st scan): >60%",
    "Ambassador referrals: >30% of signups",
    "",
    "TARGET: 6,000 users by Month 12",
], accent=CYAN)

add_card(s, Inches(9.0), Inches(1.8), Inches(3.8), Inches(4.5), "⚙️ Cost Metrics", [
    "Infra cost as % of revenue: <5%",
    "Cost per user: <$1/mo",
    "OpenAI cost per agent run: <$0.03",
    "CAC (customer acquisition): <$25",
    "Gross margin: >85%",
    "",
    "TARGET: Cash-positive by Month 6",
], accent=AMBER)


# ══════════════════════════════════════════════
# SLIDE 14: Year 1 Summary
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
   "GigLift", size=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
   "Year 1 Financial Outlook", size=28, color=GOLD, align=PP_ALIGN.CENTER)

for i, (num, label) in enumerate([
    ("$165K", "Year 1 Revenue"),
    ("$139K", "Net Income"),
    ("84%", "Net Margin"),
    ("Month 6", "Cash Positive"),
]):
    x = Inches(0.8) + Inches(3.1) * i
    tx(s, x, Inches(3.8), Inches(2.8), Inches(0.7), num, size=40, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x, Inches(4.6), Inches(2.8), Inches(0.4), label, size=14, color=MUTED, align=PP_ALIGN.CENTER)

tx(s, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
   "Lean infrastructure • Serverless architecture • <$10K startup capital • Self-funding by Month 6",
   size=16, color=LIGHT, align=PP_ALIGN.CENTER)

tx(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
   "All projections as of March 2026. Base case scenario.", size=12, color=MUTED, align=PP_ALIGN.CENTER)


# Save
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
