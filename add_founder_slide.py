#!/usr/bin/env python3
"""Add a Management Team / Founder slide to key GigLift business PPTX documents."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys

# ── Style constants (matching existing decks) ──
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(150, 150, 170)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
TEAL = RGBColor(45, 212, 191)
AMBER = RGBColor(251, 191, 36)
GOLD = RGBColor(234, 179, 8)
DARK_BG = RGBColor(10, 15, 30)
CARD_BG = RGBColor(18, 22, 42)


def add_bg(slide, W, H):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    overlay = slide.shapes.add_shape(1, Emu(0), Emu(0), W, H)
    overlay.line.fill.background()
    fill = overlay.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(10, 15, 30)
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(12, 30, 45)
    fill.gradient_stops[1].position = 1.0
    accent_line = slide.shapes.add_shape(1, Emu(0), Emu(0), W, Inches(0.03))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = TEAL
    accent_line.line.fill.background()


def tx(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(40, 50, 75)
    shape.line.width = Pt(1)
    tx(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
       title, size=15, color=accent, bold=True)
    y = top + Inches(0.55)
    for item in items:
        prefix = f"• {item}" if item and not item.startswith("•") and not item.startswith(" ") else item
        tx(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3),
           prefix, size=12, color=LIGHT)
        y += Inches(0.27)


def add_founder_slide(prs, insert_position=None):
    """Add a Management Team / Founder slide to the presentation."""
    W = prs.slide_width
    H = prs.slide_height

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, W, H)

    # Header
    tx(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
       "Management Team", size=34, color=WHITE, bold=True)
    tx(s, Inches(0.8), Inches(1.05), Inches(11), Inches(0.3),
       "Serial entrepreneur with a proven pattern of building, scaling, and exiting vertical SaaS", size=15, color=TEAL)

    # Founder card
    add_card(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.3), "BLAKE THORNLEY  —  FOUNDER & CTO", [
        "25+ years building and scaling cloud-native SaaS platforms",
        "Three-time co-founder with 3 successful exits via acquisition",
        "Led globally distributed engineering teams of 30+",
        "99.99% uptime, reduced infrastructure costs by 50%",
        "Deep expertise in cybersecurity & SOC 2 compliance",
        "",
        "Education:",
        "  Stanford University — Software Security, Cryptography, Startup Engineering",
        "  University of Redlands — Management Information Systems",
        "",
        "Languages: English (native), Spanish (fluent), French (intermediate)",
    ], accent=PURPLE)

    # Exit track record card
    add_card(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.3), "TRACK RECORD: 3 EXITS", [
        "✅  D4 Sports, Inc.  (Co-Founder / VP Technology)",
        "      League management SaaS → founding team to 25 employees",
        "      $500K+ ARR → Acquired by Clubspaces",
        "",
        "✅  Clubspaces, Inc.  (VP, IT & Architecture)",
        "      30-person engineering dept, $1M+ revenue, 1M+ users",
        "      LA Galaxy & DC United contracts → Acquired by Active Network",
        "",
        "✅  Total Global Sports  (Co-Founder / CTO / Board)",
        "      99.99% uptime, 12 engineers, 50% AWS cost reduction",
        "      8-year run → Acquired by AthleteOne",
        "",
        "Now applying the same playbook to music tech with GigLift",
    ], accent=GREEN)

    # Key stat badges at bottom
    stats = [
        ("3", "Successful Exits", TEAL),
        ("25+", "Years Experience", CYAN),
        ("30+", "Team Size Led", AMBER),
        ("99.99%", "Uptime Record", GREEN),
        ("$1M+", "Revenue Scaled", PURPLE),
    ]
    for i, (val, label, color) in enumerate(stats):
        x = Inches(0.6) + Inches(2.5) * i
        tx(s, x, Inches(7.05), Inches(2.2), Inches(0.35), val, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)

    return s


def process_document(filepath):
    """Open a PPTX, add the founder slide, and save."""
    if not os.path.exists(filepath):
        print(f"  ⚠️  Not found: {filepath}")
        return False

    try:
        prs = Presentation(filepath)
        add_founder_slide(prs)
        prs.save(filepath)
        print(f"  ✅  Updated: {os.path.basename(filepath)} ({len(prs.slides)} slides)")
        return True
    except Exception as e:
        print(f"  ❌  Failed: {os.path.basename(filepath)} — {e}")
        return False


if __name__ == "__main__":
    docs_dir = os.path.join(os.path.dirname(__file__), "data", "docs", "business")

    targets = [
        os.path.join(docs_dir, "Bandsintown_x_GigLift_MA_Deck.pptx"),
        os.path.join(docs_dir, "GigLift_Business_Plan.pptx"),
        os.path.join(docs_dir, "GigLift_Pitch_Deck.pptx"),
        os.path.join(docs_dir, "GigLift_Strategic_Outlook.pptx"),
        os.path.join(docs_dir, "GigLift_Marketing_Sales_Plan.pptx"),
        os.path.join(docs_dir, "GigLift_Brand_Ambassador_Plan.pptx"),
        os.path.join(docs_dir, "GigLift_Year1_Financial_Plan.pptx"),
    ]

    print("🏗️  Adding Management Team / Founder slide to business documents...\n")
    updated = 0
    for t in targets:
        if process_document(t):
            updated += 1

    # Also update the copies in ~/Documents if they exist
    home_docs = os.path.expanduser("~/Documents")
    home_targets = [
        os.path.join(home_docs, "GigLift_Business_Plan.pptx"),
        os.path.join(home_docs, "GigLift_Pitch_Deck.pptx"),
        os.path.join(home_docs, "GigLift_Strategic_Outlook.pptx"),
        os.path.join(home_docs, "GigLift_Marketing_Sales_Plan.pptx"),
    ]
    for t in home_targets:
        if process_document(t):
            updated += 1

    print(f"\n✅  Done! Updated {updated} documents with Management Team slide.")
