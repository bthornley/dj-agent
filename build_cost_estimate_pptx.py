#!/usr/bin/env python3
"""Build GigLift Launch Cost Estimate PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT = os.path.join(os.path.dirname(__file__), "data", "docs", "business", "GigLift_Launch_Cost_Estimate.pptx")
BG = "/Users/bthornley/.gemini/antigravity/brain/d900a025-a940-458f-9e2e-9809b998c050/ip_slide_bg_1772693744829.png"

W = Inches(13.333)
H = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(200, 200, 220)
MUTED = RGBColor(150, 150, 170)
PURPLE = RGBColor(168, 85, 247)
CYAN = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
GOLD = RGBColor(234, 179, 8)
DARK_BG = RGBColor(12, 12, 30)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    try:
        slide.shapes.add_picture(BG, Emu(0), Emu(0), W, H)
    except Exception:
        pass


def tx(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf


def add_card(slide, left, top, width, height, title, items, accent=PURPLE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(20, 20, 42)
    shape.line.color.rgb = RGBColor(50, 50, 80)
    shape.line.width = Pt(1)
    tx(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
       title, size=16, color=accent, bold=True)
    y = top + Inches(0.55)
    for item in items:
        tx(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3),
           f"• {item}", size=12, color=LIGHT)
        y += Inches(0.28)


# ══════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(1.3), Inches(11), Inches(0.6),
   "💰 LAUNCH COST ESTIMATE", size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(2.1), Inches(11), Inches(1.5),
   "Third-Party Services\n& Infrastructure Costs", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
   "Pre-Launch  •  Launch (0–1K Users)  •  Growth (1K–15K Users)", size=22, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
   "GigLift  •  March 2026", size=16, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 2: Current Stack Overview
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Technology Stack", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "8 third-party services powering GigLift", size=16, color=GOLD)

services = [
    ("⚡ Vercel", CYAN, ["Next.js hosting & serverless", "Edge functions & cron jobs", "Blob storage for media"]),
    ("🔐 Clerk", PURPLE, ["Authentication & user mgmt", "Social logins, MFA", "Role-based access control"]),
    ("🗄️ Turso", GREEN, ["Edge SQLite (libSQL)", "Globally distributed", "All application data"]),
    ("💳 Stripe", AMBER, ["Payment processing", "Subscription management", "Billing portal"]),
]

for i, (name, color, items) in enumerate(services):
    x = Inches(0.5) + Inches(3.1) * i
    add_card(s, x, Inches(1.8), Inches(2.9), Inches(2.8), name, items, accent=color)

services2 = [
    ("📧 Resend", CYAN, ["Transactional emails", "Outreach & notifications"]),
    ("🤖 OpenAI", PURPLE, ["GPT-4o-mini for agents", "Lead scoring, content gen"]),
    ("📊 Plausible", GREEN, ["Privacy-first analytics", "No cookies required"]),
    ("🌐 Domains", AMBER, ["giglift.com (GoDaddy)", "giglift.app (Google)"]),
]

for i, (name, color, items) in enumerate(services2):
    x = Inches(0.5) + Inches(3.1) * i
    add_card(s, x, Inches(4.9), Inches(2.9), Inches(2.3), name, items, accent=color)

# ══════════════════════════════════════════════
# SLIDE 3: Pre-Launch (Current Costs)
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Phase 1: Pre-Launch (Now)", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Current monthly costs — most services on free tiers", size=16, color=GOLD)

# Header
tx(s, Inches(0.8), Inches(1.9), Inches(4), Inches(0.3), "SERVICE", size=12, color=MUTED, bold=True)
tx(s, Inches(5.2), Inches(1.9), Inches(2.2), Inches(0.3), "PLAN", size=12, color=MUTED, bold=True)
tx(s, Inches(7.6), Inches(1.9), Inches(1.8), Inches(0.3), "MONTHLY", size=12, color=MUTED, bold=True)
tx(s, Inches(9.6), Inches(1.9), Inches(2.5), Inches(0.3), "NOTES", size=12, color=MUTED, bold=True)

items = [
    ("Vercel", "Hobby", "$0", "Free for personal/non-commercial", GREEN),
    ("Clerk", "Free", "$0", "Up to 50K MAU included", GREEN),
    ("Turso (libSQL)", "Free", "$0", "500 DBs, 9 GB storage, 1B reads", GREEN),
    ("Stripe", "Pay-as-you-go", "$0", "No monthly fee, 2.9% + 30¢/txn", GREEN),
    ("Resend", "Free", "$0", "3K emails/mo, 100/day limit", GREEN),
    ("OpenAI API", "Pay-as-you-go", "~$5–15", "7 agents × daily cron runs", AMBER),
    ("Plausible", "Growth", "$9", "Up to 10K monthly pageviews", AMBER),
    ("Domains", "—", "~$3", "giglift.com + giglift.app (~$40/yr)", AMBER),
]

y = Inches(2.3)
for name, plan, monthly, notes, color in items:
    tx(s, Inches(0.8), y, Inches(4), Inches(0.3), name, size=14, color=LIGHT)
    tx(s, Inches(5.2), y, Inches(2.2), Inches(0.3), plan, size=14, color=WHITE)
    tx(s, Inches(7.6), y, Inches(1.8), Inches(0.3), monthly, size=14, color=color, bold=True)
    tx(s, Inches(9.6), y, Inches(3), Inches(0.3), notes, size=11, color=MUTED)
    y += Inches(0.45)

tx(s, Inches(0.8), y + Inches(0.2), Inches(4), Inches(0.35), "TOTAL (PRE-LAUNCH)", size=16, color=GOLD, bold=True)
tx(s, Inches(7.6), y + Inches(0.2), Inches(2), Inches(0.35), "~$17–27/mo", size=16, color=GOLD, bold=True)

# ══════════════════════════════════════════════
# SLIDE 4: Launch Phase (0–1K Users)
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Phase 2: Launch (0–1K Users)", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Upgrade to production plans — Vercel Pro required for commercial use", size=16, color=GOLD)

tx(s, Inches(0.8), Inches(1.9), Inches(4), Inches(0.3), "SERVICE", size=12, color=MUTED, bold=True)
tx(s, Inches(5.2), Inches(1.9), Inches(2.2), Inches(0.3), "PLAN", size=12, color=MUTED, bold=True)
tx(s, Inches(7.6), Inches(1.9), Inches(1.8), Inches(0.3), "MONTHLY", size=12, color=MUTED, bold=True)
tx(s, Inches(9.6), Inches(1.9), Inches(2.5), Inches(0.3), "NOTES", size=12, color=MUTED, bold=True)

items = [
    ("Vercel", "Pro", "$20", "Required for commercial use", AMBER),
    ("Clerk", "Free → Pro", "$0–20", "Free to 50K MAU; $20/mo removes branding", GREEN),
    ("Turso", "Scaler", "$29", "24 GB, 100B reads, PITR 30 days", AMBER),
    ("Stripe", "Pay-as-you-go", "$0 + fees", "2.9% + 30¢ per transaction", GREEN),
    ("Resend", "Pro", "$20", "50K emails/mo, no daily limit", AMBER),
    ("OpenAI API", "Pay-as-you-go", "~$20–50", "Higher agent volume as users grow", AMBER),
    ("Plausible", "Growth", "$19", "Up to 100K monthly pageviews", AMBER),
    ("Domains", "—", "~$3", "Annual renewals", GREEN),
    ("Vercel Blob", "Included", "$0–5", "$0.023/GB, 5 GB free on Pro", GREEN),
]

y = Inches(2.3)
for name, plan, monthly, notes, color in items:
    tx(s, Inches(0.8), y, Inches(4), Inches(0.3), name, size=14, color=LIGHT)
    tx(s, Inches(5.2), y, Inches(2.2), Inches(0.3), plan, size=14, color=WHITE)
    tx(s, Inches(7.6), y, Inches(1.8), Inches(0.3), monthly, size=14, color=color, bold=True)
    tx(s, Inches(9.6), y, Inches(3), Inches(0.3), notes, size=11, color=MUTED)
    y += Inches(0.45)

tx(s, Inches(0.8), y + Inches(0.2), Inches(4), Inches(0.35), "TOTAL (LAUNCH)", size=16, color=GOLD, bold=True)
tx(s, Inches(7.6), y + Inches(0.2), Inches(2), Inches(0.35), "~$111–166/mo", size=16, color=GOLD, bold=True)

# ══════════════════════════════════════════════
# SLIDE 5: Growth Phase (1K–15K Users)
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Phase 3: Growth (1K–15K Users)", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Year 1–2 scaling — costs grow proportionally with revenue", size=16, color=GOLD)

tx(s, Inches(0.8), Inches(1.9), Inches(4), Inches(0.3), "SERVICE", size=12, color=MUTED, bold=True)
tx(s, Inches(5.2), Inches(1.9), Inches(2.2), Inches(0.3), "PLAN", size=12, color=MUTED, bold=True)
tx(s, Inches(7.6), Inches(1.9), Inches(1.8), Inches(0.3), "MONTHLY", size=12, color=MUTED, bold=True)
tx(s, Inches(9.6), Inches(1.9), Inches(2.5), Inches(0.3), "NOTES", size=12, color=MUTED, bold=True)

items = [
    ("Vercel", "Pro", "$20–60", "Usage-based overage beyond $20 credit", AMBER),
    ("Clerk", "Pro", "$20 + MAU", "$0.02/MAU beyond 50K", AMBER),
    ("Turso", "Scaler → Pro", "$29–499", "Scale as DB grows; Pro at $499/mo", RED),
    ("Stripe", "Pay-as-you-go", "$0 + fees", "~$870–$5,100/mo at scale (2.9%+30¢)", AMBER),
    ("Resend", "Scale", "$90", "100K emails/mo", AMBER),
    ("OpenAI API", "Pay-as-you-go", "~$100–300", "Per-user agent runs scale linearly", RED),
    ("Plausible", "Business", "$29–39", "Up to 1M monthly pageviews", AMBER),
    ("Domains", "—", "~$3", "Annual renewals", GREEN),
    ("Vercel Blob", "Pro", "$5–20", "Media uploads grow with users", AMBER),
]

y = Inches(2.3)
for name, plan, monthly, notes, color in items:
    tx(s, Inches(0.8), y, Inches(4), Inches(0.3), name, size=14, color=LIGHT)
    tx(s, Inches(5.2), y, Inches(2.2), Inches(0.3), plan, size=14, color=WHITE)
    tx(s, Inches(7.6), y, Inches(1.8), Inches(0.3), monthly, size=14, color=color, bold=True)
    tx(s, Inches(9.6), y, Inches(3), Inches(0.3), notes, size=11, color=MUTED)
    y += Inches(0.45)

tx(s, Inches(0.8), y + Inches(0.2), Inches(4), Inches(0.35), "TOTAL (GROWTH)", size=16, color=GOLD, bold=True)
tx(s, Inches(7.6), y + Inches(0.2), Inches(2), Inches(0.35), "~$296–$1,000/mo", size=16, color=GOLD, bold=True)

# ══════════════════════════════════════════════
# SLIDE 6: OpenAI Deep Dive
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "OpenAI API — Largest Variable Cost", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "7 autonomous agents run daily via cron — costs scale with user count", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5), "Token Pricing (GPT-4o-mini)", [
    "Input: $0.15 / 1M tokens",
    "Output: $0.60 / 1M tokens",
    "Cached input: $0.075 / 1M tokens",
    "",
    "Each agent run: ~2K–5K tokens",
    "7 agents × 365 days = 2,555 runs/yr (system-level)",
    "Per-user agents (lead scan, social): ~$0.01–0.03/run",
], accent=PURPLE)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.5), "Estimated Monthly Cost by Scale", [
    "Pre-launch (system agents only): ~$5–15",
    "100 active users: ~$20–40",
    "1,000 active users: ~$50–100",
    "5,000 active users: ~$150–250",
    "15,000 active users: ~$300–500",
    "",
    "GPT-4o (if needed): 17x more expensive",
], accent=RED)

tx(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2),
   "💡 Optimization levers: Use cached prompts for repeated patterns, batch agent runs during "
   "off-peak hours, throttle per-user agent frequency on free tier, consider open-source models "
   "(Llama, Mistral) for lower-complexity tasks to reduce costs 80%+.", size=13, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 7: Stripe Revenue vs. Cost
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Stripe Fees vs. Revenue", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Payment processing is purely variable — 2.9% + $0.30 per transaction", size=16, color=GOLD)

tx(s, Inches(0.8), Inches(1.9), Inches(3), Inches(0.3), "SUBSCRIBERS", size=12, color=MUTED, bold=True)
tx(s, Inches(3.8), Inches(1.9), Inches(2), Inches(0.3), "AVG PLAN", size=12, color=MUTED, bold=True)
tx(s, Inches(5.8), Inches(1.9), Inches(2.2), Inches(0.3), "MONTHLY REV", size=12, color=MUTED, bold=True)
tx(s, Inches(8.2), Inches(1.9), Inches(2), Inches(0.3), "STRIPE FEES", size=12, color=MUTED, bold=True)
tx(s, Inches(10.2), Inches(1.9), Inches(2), Inches(0.3), "NET", size=12, color=MUTED, bold=True)

scenarios = [
    ("100", "$33/mo", "$3,300", "~$126", "$3,174"),
    ("500", "$33/mo", "$16,500", "~$629", "$15,871"),
    ("2,500", "$33/mo", "$82,500", "~$3,143", "$79,357"),
    ("5,000", "$38/mo", "$190,000", "~$7,010", "$182,990"),
    ("15,000", "$38/mo", "$570,000", "~$21,030", "$548,970"),
]

y = Inches(2.3)
for subs, plan, rev, fees, net in scenarios:
    tx(s, Inches(0.8), y, Inches(3), Inches(0.3), subs, size=14, color=WHITE, bold=True)
    tx(s, Inches(3.8), y, Inches(2), Inches(0.3), plan, size=14, color=LIGHT)
    tx(s, Inches(5.8), y, Inches(2.2), Inches(0.3), rev, size=14, color=GREEN)
    tx(s, Inches(8.2), y, Inches(2), Inches(0.3), fees, size=14, color=RED)
    tx(s, Inches(10.2), y, Inches(2), Inches(0.3), net, size=14, color=GREEN, bold=True)
    y += Inches(0.48)

tx(s, Inches(0.8), y + Inches(0.3), Inches(11), Inches(0.5),
   "Stripe has no monthly fee. International cards: +1.5%. Stripe Tax add-on: 0.5% per txn if needed.", size=12, color=MUTED)

# ══════════════════════════════════════════════
# SLIDE 8: One-Time Launch Costs
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "One-Time Launch Costs", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Non-recurring costs required before or at launch", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), "Required", [
    "Domain portfolio (giglift.io, .co, .ai): ~$50–150",
    "Clerk custom domain setup: $0 (included on Pro)",
    "Stripe account verification: $0",
    "Resend domain verification (DNS): $0",
    "Turso production DB setup: $0",
    "Apple Developer Program (if iOS app): $99/yr",
    "Google Play Console (if Android): $25 one-time",
    "",
    "SSL certificates: $0 (auto via Vercel)",
    "CI/CD: $0 (included with Vercel)",
], accent=GREEN)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), "Recommended", [
    "Trademark registration (2 classes): $2,000–3,000",
    "Provisional patent: $2,100–5,100",
    "Terms of Service / Privacy Policy (attorney): $1,000–3,000",
    "SOC 2 compliance prep (if enterprise): $5,000–15,000",
    "Penetration testing: $2,000–5,000",
    "Logo / brand design (if refreshing): $500–2,000",
    "",
    "Total recommended: $12,600–$33,100",
], accent=AMBER)

# ══════════════════════════════════════════════
# SLIDE 9: Annual Cost Summary
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Annual Cost Summary", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Total cost of third-party services by phase", size=16, color=GOLD)

phases = [
    ("Pre-Launch", AMBER, [
        "Monthly: ~$17–27",
        "Annual: ~$200–320",
        "",
        "Mostly free tiers",
        "Only OpenAI + Plausible + domains",
    ]),
    ("Launch (0–1K)", CYAN, [
        "Monthly: ~$111–166",
        "Annual: ~$1,330–2,000",
        "",
        "Vercel Pro + Turso Scaler + Resend Pro",
        "Still within free Clerk tier",
    ]),
    ("Growth (1K–15K)", PURPLE, [
        "Monthly: ~$296–1,000",
        "Annual: ~$3,550–12,000",
        "",
        "OpenAI becomes largest variable cost",
        "Turso may need Pro ($499/mo) at scale",
    ]),
]

for i, (name, color, items) in enumerate(phases):
    x = Inches(0.6) + Inches(4.1) * i
    add_card(s, x, Inches(1.8), Inches(3.8), Inches(4.0), name, items, accent=color)

tx(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8),
   "Infrastructure cost as % of revenue: ~3–5% at launch, dropping to ~1–2% at scale. "
   "This is well within SaaS industry benchmarks (typically 5–15% of revenue).", size=14, color=LIGHT)

# ══════════════════════════════════════════════
# SLIDE 10: Cost Risks & Optimization
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
   "Cost Risks & Optimizations", size=36, color=WHITE, bold=True)
tx(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Key risks and strategies to keep costs manageable", size=16, color=GOLD)

add_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), "⚠️ Cost Risks", [
    "OpenAI price increases or rate limit changes",
    "Turso jump from Scaler ($29) to Pro ($499)",
    "Clerk MAU overage at scale ($0.02/user)",
    "Vercel serverless function timeout limits",
    "Stripe international card surcharges",
    "Unexpected DDoS or bot traffic spikes",
    "GDPR compliance costs (EU expansion)",
    "",
    "Worst-case Y1 monthly: ~$1,500–2,000",
], accent=RED)

add_card(s, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), "✅ Optimization Strategies", [
    "Use GPT-4o-mini (17x cheaper than GPT-4o)",
    "Cache repeated AI prompts (50% cost reduction)",
    "Rate-limit free-tier users on agent runs",
    "Batch agent cron jobs during off-peak hours",
    "Monitor Turso row reads to stay on Scaler",
    "Use Vercel ISR/static pages where possible",
    "Consider self-hosted Llama/Mistral for 80% savings",
    "",
    "Target: keep infra <5% of ARR at all times",
], accent=GREEN)

# ══════════════════════════════════════════════
# SLIDE 11: Closing
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tx(s, Inches(1), Inches(2), Inches(11), Inches(1),
   "GigLift", size=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
   "Launch Cost Estimate", size=28, color=GOLD, align=PP_ALIGN.CENTER)

for i, (num, label) in enumerate([
    ("~$17/mo", "Pre-Launch"),
    ("~$140/mo", "Launch"),
    ("~$650/mo", "Growth"),
    ("<3%", "Infra % of ARR"),
]):
    x = Inches(0.8) + Inches(3.1) * i
    tx(s, x, Inches(4.3), Inches(2.8), Inches(0.7), num, size=36, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x, Inches(5.0), Inches(2.8), Inches(0.4), label, size=14, color=MUTED, align=PP_ALIGN.CENTER)

tx(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
   "All pricing as of March 2026. Verify current rates before committing.", size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Save
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   {len(prs.slides)} slides")
